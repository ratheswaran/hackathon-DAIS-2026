#!/usr/bin/env python3
"""Build ``ra_template.pptx`` — the Resonance Analytics deck template that turns
``compose_deck`` template mode ON.

Why this exists
---------------
``compose_deck`` (``hackathon-orchestrator/tools/compose_deck.py``) renders decks
from a JSON ``deck_spec``. By default it paints the full RA brand chrome itself
(*from-scratch mode*). When a ``ra_template.pptx`` sits in this folder the tool
switches to *template mode*: ``_TEMPLATE_MODE`` early-returns the chrome helpers
(``_add_background``, ``_draw_brandmark``, ``_draw_cover_orb``,
``_draw_slide_foot``) and paints **only content** on top of named slide layouts.
Backgrounds, brand-mark, wordmark, footer, page number and the cover orb then
come from THIS template's master + layouts.

The renderer's exact contract (see ``_render_pptx`` in ``compose_deck.py``):

* It loads the template, keeps only **slide 0** as the cover seed; for a
  ``cover`` first slide it fills that seed's placeholder **idx 10** (title,
  Source Serif) + **idx 1** (subtitle).
* A ``closing`` slide is ``add_slide(RA_thankyou)`` → fills **idx 11** (title) +
  **idx 1** (subtitle).
* Every other slide is ``add_slide(RA_blank | RA_dark | RA_SectionSignal |
  RA_SectionWhite)``; empty non-title placeholders are stripped, then content is
  painted on top — so content/section layouts must carry NO content
  placeholders, just background + chrome.

So this builder produces six layouts named exactly:
    RA_TitleSlide · RA_blank · RA_dark · RA_SectionSignal · RA_SectionWhite · RA_thankyou
with the cover/closing titles re-indexed to 10/11 + a subtitle at idx 1, and the
stock Date(10)/Footer(11)/SlideNumber(12) placeholders deleted (they'd collide
with the renderer's title indices).

Chrome comes from the PNG assets in ``../assets`` (cover poster, wordmarks,
geo-mark). Because python-pptx's high-level ``add_picture``/``add_textbox`` are
slide-only, layout chrome is injected as raw OOXML onto the layout part (images
embedded via ``part.get_or_add_image_part``). Page numbers use a live
``slidenum`` FIELD in a plain textbox (a slide-number *placeholder* would be
stripped by the renderer; a field in a normal textbox survives).

Run:  python3 build_ra_template.py        # writes ra_template.pptx next to this
"""

from __future__ import annotations

import os

from lxml import etree
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "..", "assets")
OUT = os.path.join(HERE, "ra_template.pptx")

W_IN, H_IN = 13.333, 7.5            # 16:9 widescreen — matches SLIDE_W_IN/H_IN
EMU = 914400

# Resonance Analytics palette (sRGB renderings of the oklch tokens).
SIGNAL_600, SIGNAL_700 = "254BB2", "17368D"
INK_950, INK_900, INK_500, INK_300 = "010204", "050A0F", "6D7277", "C6CBD0"
PAPER, CREAM, GOLD = "FFFFFF", "FBF0E4", "DF9B44"
# Subtle cool-blue paper = the design system's real --paper (--ink-50,
# oklch(0.985 0.003 250)). Used ONLY for RA_SectionWhite per the brand call.
PAPER_TINT = "F9FAFC"

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Intrinsic asset pixel dims (avoids a Pillow dependency at build time).
IMG_WH = {
    "wordmark-dark.png": (900, 128), "wordmark-light.png": (900, 128),
    # Official RA logo lockups: the "reso" radial mark + "Resonance Analytics"
    # in Manrope (dark = white text for dark bg, light = ink text for light bg).
    "ra-lockup-dark.png": (2733, 320), "ra-lockup-light.png": (2733, 320),
    "ra-logo-dark.png": (643, 640), "ra-logo-light.png": (643, 640),
    "geomark-amber.png": (480, 480), "brandmark-signal.png": (200, 200),
    "brandmark-white.png": (200, 200), "cover-poster.png": (1920, 1080),
    "cover-orb.png": (1400, 1400),
}

_ID = [900]


def _nid() -> int:
    _ID[0] += 1
    return _ID[0]


def emu(inches: float) -> int:
    return int(round(inches * EMU))


def asset(name: str) -> str:
    return os.path.join(ASSETS, name)


# --------------------------------------------------------------------------- #
# Raw-OOXML shape injectors — work on any host (slide / layout / master)
# --------------------------------------------------------------------------- #
_PIC_XML = (
    '<p:pic xmlns:p="%s" xmlns:a="%s" xmlns:r="%s">'
    '<p:nvPicPr><p:cNvPr id="%d" name="%s"/>'
    '<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
    '<p:blipFill><a:blip r:embed="%s"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
    '<p:spPr><a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/></a:xfrm>'
    '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
)

_FLD_SP = (
    '<p:sp xmlns:p="%s" xmlns:a="%s">'
    '<p:nvSpPr><p:cNvPr id="%d" name="%s"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
    '<p:spPr><a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/></a:xfrm>'
    '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
    '<p:txBody>'
    '<a:bodyPr wrap="none" lIns="0" tIns="0" rIns="0" bIns="0" anchor="ctr"><a:noAutofit/></a:bodyPr>'
    '<a:lstStyle/><a:p><a:pPr algn="%s"/>'
    '<a:fld id="{B7C27A6A-2B1E-4F2C-9E11-1A2B3C4D5E6F}" type="slidenum">'
    '<a:rPr lang="en-US" sz="%d" b="0"><a:solidFill><a:srgbClr val="%s"/></a:solidFill>'
    '<a:latin typeface="JetBrains Mono"/></a:rPr><a:t>2</a:t></a:fld>'
    '</a:p></p:txBody></p:sp>'
)


def add_pic(host, fname: str, x: float, y: float, *, w=None, h=None):
    """Embed + place an image on any slide/layout host. Returns the <p:pic> el."""
    iw, ih = IMG_WH[fname]
    if w is None and h is not None:
        w = h * iw / ih
    elif h is None and w is not None:
        h = w * ih / iw
    elif w is None and h is None:
        w, h = iw / 96.0, ih / 96.0
    _, rId = host.part.get_or_add_image_part(asset(fname))
    el = parse_xml(_PIC_XML % (P_NS, A_NS, R_NS, _nid(), fname.split(".")[0],
                               rId, emu(x), emu(y), emu(w), emu(h)))
    host.shapes._spTree.append(el)
    return el


def add_pagenum(host, x: float, y: float, w: float, h: float, *, color: str,
                size: float = 10.0, align: str = "r"):
    el = parse_xml(_FLD_SP % (P_NS, A_NS, _nid(), "pagenum",
                              emu(x), emu(y), emu(w), emu(h),
                              align, int(size * 100), color))
    host.shapes._spTree.append(el)
    return el


def add_footer(host, *, dark: bool) -> None:
    """Footer lockup: small RA logo lockup left + page-number field right."""
    wm = "ra-lockup-dark.png" if dark else "ra-lockup-light.png"  # dark=white text
    add_pic(host, wm, 0.75, 7.04, h=0.18)
    add_pagenum(host, 11.0, 6.96, 1.583, 0.34,
                color=(INK_300 if dark else INK_500), size=10)


# --------------------------------------------------------------------------- #
# Placeholder + background helpers
# --------------------------------------------------------------------------- #
def set_layout_name(layout, name: str) -> None:
    layout._element.cSld.set("name", name)


def reindex_ph(shape, idx: int, ph_type: str = "body") -> None:
    el = shape._element.nvSpPr.nvPr.find(qn("p:ph"))
    el.set("idx", str(idx))
    if ph_type:
        el.set("type", ph_type)
    elif "type" in el.attrib:
        del el.attrib["type"]


def delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def delete_all_placeholders(layout) -> None:
    for ph in list(layout.placeholders):
        delete_shape(ph)


def place(shape, x, y, w, h) -> None:
    shape.left, shape.top, shape.width, shape.height = (
        Inches(x), Inches(y), Inches(w), Inches(h))


def set_bg(host, hex_color: str) -> None:
    """Insert a solid-fill <p:bg> as the first child of <p:cSld>."""
    cSld = host._element.cSld
    for existing in cSld.findall(qn("p:bg")):
        cSld.remove(existing)
    bg = parse_xml(
        '<p:bg xmlns:p="%s" xmlns:a="%s"><p:bgPr><a:solidFill>'
        '<a:srgbClr val="%s"/></a:solidFill><a:effectLst/></p:bgPr></p:bg>'
        % (P_NS, A_NS, hex_color))
    cSld.insert(0, bg)


def style_placeholder(shape, *, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT) -> None:
    """Left/top, zero inset, so the renderer's run lands where the box sits."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    txBody = shape._element.find(qn("p:txBody"))
    if txBody is None:
        return
    for old in txBody.findall(qn("a:lstStyle")):
        txBody.remove(old)
    algn = {PP_ALIGN.LEFT: "l", PP_ALIGN.CENTER: "ctr", PP_ALIGN.RIGHT: "r"}.get(align, "l")
    # marL/indent 0 + buNone: kill the inherited body bullet + hanging indent so
    # the renderer's title/subtitle run starts flush-left with no "• " glyph.
    lst = parse_xml('<a:lstStyle xmlns:a="%s"><a:lvl1pPr marL="0" indent="0" algn="%s">'
                    '<a:buNone/><a:defRPr/></a:lvl1pPr></a:lstStyle>' % (A_NS, algn))
    txBody.find(qn("a:bodyPr")).addnext(lst)


# --------------------------------------------------------------------------- #
# Theme — brand fonts + accent (best effort; runs set explicit fonts anyway)
# --------------------------------------------------------------------------- #
def set_theme(prs) -> None:
    try:
        tp = prs.slide_masters[0].part.part_related_by(RT.THEME)
        root = etree.fromstring(tp.blob)
        te = root.find(qn("a:themeElements"))
        fs = te.find(qn("a:fontScheme"))
        fs.find(qn("a:majorFont")).find(qn("a:latin")).set("typeface", "Playfair Display")
        fs.find(qn("a:minorFont")).find(qn("a:latin")).set("typeface", "Manrope")
        acc1 = te.find(qn("a:clrScheme")).find(qn("a:accent1"))
        for child in list(acc1):
            acc1.remove(child)
        acc1.append(parse_xml('<a:srgbClr xmlns:a="%s" val="%s"/>' % (A_NS, SIGNAL_600)))
        tp._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8",
                                  standalone=True)
    except Exception as e:  # noqa: BLE001
        print("  ! theme tweak skipped:", e)


# --------------------------------------------------------------------------- #
# Build
# --------------------------------------------------------------------------- #
def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    set_theme(prs)

    L = list(prs.slide_layouts)         # 11 stock layouts; keep 6, drop 5
    master = prs.slide_masters[0]

    # ---- RA_TitleSlide  (stock 0 "Title Slide": ctrTitle + subTitle) -------
    title_lay = L[0]
    set_layout_name(title_lay, "RA_TitleSlide")
    set_bg(title_lay, INK_950)
    phs = {ph.placeholder_format.idx: ph for ph in title_lay.placeholders}
    for idx in (10, 11, 12):            # delete stock Date/Footer/SlideNumber
        if idx in phs:
            delete_shape(phs[idx])
    reindex_ph(phs[0], 10, "body")      # CENTER_TITLE -> idx 10
    place(phs[0], 0.75, 2.85, 11.83, 2.3)
    style_placeholder(phs[0])
    reindex_ph(phs[1], 1, "body")       # SUBTITLE     -> idx 1
    place(phs[1], 0.78, 5.25, 10.5, 1.0)
    style_placeholder(phs[1])
    add_pic(title_lay, "ra-lockup-dark.png", 0.75, 0.52, h=0.27)

    # ---- RA_blank  (stock 6 "Blank": no content placeholders) --------------
    blank_lay = L[6]
    set_layout_name(blank_lay, "RA_blank")
    set_bg(blank_lay, PAPER)
    delete_all_placeholders(blank_lay)
    add_footer(blank_lay, dark=False)

    # ---- RA_cream  (stock 1 "Title and Content" -> stripped: warm blank) ---
    # Content layout for bg:"cream" slides (quotes, data-notes). Mirrors
    # RA_blank but on the warm brand cream — footer only, no top-left lockup,
    # so painted content never collides with a header mark.
    cream_lay = L[1]
    set_layout_name(cream_lay, "RA_cream")
    set_bg(cream_lay, CREAM)
    delete_all_placeholders(cream_lay)
    add_footer(cream_lay, dark=False)

    # ---- RA_dark  (stock 5 "Title Only" -> stripped: dark blank) -----------
    dark_lay = L[5]
    set_layout_name(dark_lay, "RA_dark")
    set_bg(dark_lay, INK_950)
    delete_all_placeholders(dark_lay)
    add_footer(dark_lay, dark=True)

    # ---- RA_SectionSignal  (stock 2 "Section Header" -> stripped) ----------
    sig_lay = L[2]
    set_layout_name(sig_lay, "RA_SectionSignal")
    set_bg(sig_lay, SIGNAL_700)
    delete_all_placeholders(sig_lay)
    add_pic(sig_lay, "ra-lockup-dark.png", 0.75, 0.52, h=0.27)
    add_footer(sig_lay, dark=True)

    # ---- RA_SectionWhite  (stock 3 "Two Content" -> stripped) --------------
    white_lay = L[3]
    set_layout_name(white_lay, "RA_SectionWhite")
    set_bg(white_lay, PAPER_TINT)        # subtle cool-blue paper, not pure white
    delete_all_placeholders(white_lay)
    add_pic(white_lay, "ra-lockup-light.png", 0.75, 0.52, h=0.27)
    add_footer(white_lay, dark=False)

    # ---- RA_thankyou  (stock 4 "Comparison": keep 2 phs as 11 + 1) ---------
    ty_lay = L[4]
    set_layout_name(ty_lay, "RA_thankyou")
    set_bg(ty_lay, INK_950)
    phs = {ph.placeholder_format.idx: ph for ph in ty_lay.placeholders}
    for idx in (2, 3, 4, 10, 11, 12):   # drop extra bodies + Date/Footer/SlideNum
        if idx in phs:
            delete_shape(phs[idx])
    reindex_ph(phs[0], 11, "body")      # TITLE -> idx 11
    place(phs[0], 0.75, 2.85, 11.83, 2.1)
    style_placeholder(phs[0])
    reindex_ph(phs[1], 1, "body")       # BODY  -> idx 1
    place(phs[1], 0.78, 5.0, 10.0, 0.9)
    style_placeholder(phs[1])
    add_pic(ty_lay, "ra-lockup-dark.png", 0.75, 0.52, h=0.27)
    add_pic(ty_lay, "geomark-amber.png", 11.7, 5.78, h=1.05)
    add_pagenum(ty_lay, 11.0, 6.96, 1.583, 0.34, color=INK_300, size=10)

    # ---- delete the 4 unused stock layouts (kept L[1] for RA_cream) --------
    for i in (7, 8, 9, 10):
        try:
            master.slide_layouts.remove(L[i])
        except Exception as e:  # noqa: BLE001
            print("  ! could not remove stock layout %d: %s" % (i, e))

    # ---- cover SEED slide (slide 0) ----------------------------------------
    # The renderer reuses slides[0] as the cover and only fills idx 10 + idx 1.
    # Poster + wordmark live ON THE SLIDE so they always render.
    seed = prs.slides.add_slide(title_lay)
    set_bg(seed, INK_950)
    poster = add_pic(seed, "cover-poster.png", 0, 0, w=W_IN, h=H_IN)
    spTree = poster.getparent()                 # send poster to the back
    spTree.remove(poster)
    spTree.insert(2, poster)
    add_pic(seed, "ra-lockup-dark.png", 0.75, 0.52, h=0.27)
    add_pagenum(seed, 11.0, 6.96, 1.583, 0.34, color=INK_300, size=10)

    prs.save(OUT)
    print("wrote", OUT)

    # Embed the brand fonts (Source Serif 4 / Manrope / JetBrains Mono) so the
    # serif + mono render on machines without them installed. Skipped quietly
    # if the bundled TTFs aren't present.
    try:
        from embed_fonts import embed_fonts, FONTS_DIR
        if os.path.isdir(FONTS_DIR):
            info = embed_fonts(OUT, OUT)
            print("embedded %d font files for %s" % (info["fonts"], info["families"]))
        else:
            print("  ! fonts dir missing, skipped embedding:", FONTS_DIR)
    except Exception as e:  # noqa: BLE001
        print("  ! font embedding skipped:", e)


if __name__ == "__main__":
    main()
