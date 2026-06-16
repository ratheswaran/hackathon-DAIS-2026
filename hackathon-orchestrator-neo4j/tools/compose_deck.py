"""Compose a presentation deck from a structured slide spec — Resonance Analytics brand.

Ported from the prior ``compose_deck`` (the 2026-05-20 JSON-spec rewrite)
and re-skinned to the **Resonance Analytics** "Boardroom Editorial" design
system (cobalt-indigo *Signal* primary, cool-slate *ink* neutrals,
amber/cyan/magenta accents, Manrope + Playfair Display). The design source of
truth is ``/skills/compose-pptx/design_system_deck_example.html`` +
``design.md``.

The agent emits a JSON ``deck_spec`` (list of slide dicts; see
``/skills/compose-pptx/spec_reference.md``) and this tool:

1. **Renders PPTX directly via python-pptx** — native PowerPoint shapes,
   text boxes, tables, and charts. Charts are real editable PowerPoint
   chart objects (bars/lines/pies), not rasterised screenshots.

2. **Renders an HTML preview from the same spec** — served by the app's
   ``/api/decks/:id`` route for in-browser preview before download.

3. **Uploads ``.pptx``, ``.html``, and the raw ``.json`` spec** to
   ``/Volumes/workspace/ai_ops/agent_scratch/documents/<id>__<slug>.*``
   (the hackathon free-edition catalog) via the Databricks Files API.

4. **Returns two URLs** for the user: ``preview_url`` (HTML) and
   ``pptx_url`` (download). No PDF path — export to PDF from PowerPoint.

Brand notes
-----------
* **No raster logos.** The RA brand-mark (a rounded *Signal* square + a
  small cross-hair) and the "Resonance Analytics" wordmark are drawn
  natively as PPTX shapes — nothing to ship as a PNG. The cover gets a
  cobalt *orb* (the design's signature) drawn from layered ovals.
* **From-scratch by default.** No ``ra_template.pptx`` is required; the
  renderer paints the full RA brand chrome itself. Drop a template at
  ``skills/compose-pptx/templates/ra_template.pptx`` (or set
  ``COMPOSE_DECK_TEMPLATE``) to switch to template mode — see the
  ``RA_*`` layout-name constants below and ``templates/README.md``.
* **Fonts.** Manrope + Playfair Display are referenced by name. The HTML
  preview loads them from Google Fonts (always faithful). In PowerPoint
  they render pixel-faithfully *if installed*; otherwise PowerPoint
  substitutes, preserving the sans/serif contrast.
"""

from __future__ import annotations

import base64
import contextvars
import hashlib
import html
import io
import json
import logging
import os
import re
from pathlib import Path
from typing import Annotated, Any

from langchain_core.runnables import RunnableConfig
from langchain_core.tools import tool
from langgraph.prebuilt import InjectedStore
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from tools.compact_ref import _compact_error

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Constants — Resonance Analytics brand palette + slide canvas
# ---------------------------------------------------------------------------

# Hackathon free-edition catalog. Matches compose_document /
# compose_infographic so all three land in the same documents folder.
_VOLUME_ROOT = "/Volumes/workspace/ai_ops/agent_scratch/documents"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Resonance Analytics tokens. Hex values are the sRGB renderings of the
# design system's oklch tokens (see tokens.css / design.md). "Signal" is the
# cobalt-indigo brand scale; "ink" is the cool-slate neutral scale.
RA = {
    # Signal — cobalt-indigo brand scale
    "signal_50":   "F0F5FF",
    "signal_100":  "DAE9FF",
    "signal_200":  "B5D2FF",
    "signal_300":  "85B0FF",
    "signal_400":  "5080EB",
    "signal_500":  "345FCF",
    "signal_600":  "254BB2",   # PRIMARY
    "signal_700":  "17368D",
    "signal_800":  "0A2163",
    "signal_900":  "051139",
    # Ink — cool-slate neutrals
    "ink_50":      "F9FAFC",
    "ink_100":     "EFF2F5",
    "ink_200":     "E1E5E9",
    "ink_300":     "C6CBD0",
    "ink_400":     "9A9FA5",
    "ink_500":     "6D7277",
    "ink_600":     "494E52",
    "ink_700":     "2F3338",
    "ink_800":     "161B20",
    "ink_900":     "050A0F",
    "ink_950":     "010204",
    # Accents — harmonized with the cobalt primary
    "gold":        "DF9B44",   # amber
    "gold_deep":   "B46D10",
    "amber_light": "ECA851",
    "teal":        "2695AC",   # cyan
    "teal_deep":   "006A84",
    "plum":        "913F82",   # magenta
    "plum_deep":   "691F5D",
    "sand":        "BDD5E6",
    "slate":       "52657A",
    # Surfaces
    "white":       "FFFFFF",
    "paper":       "FFFFFF",
    "cream":       "FBF0E4",   # warm off-white
    "cover_base":  "01030B",   # near-black cobalt-tinted (cover poster base)
    # Status
    "success":     "14874E",
    "warning":     "D79628",
    "danger":      "BA2B2E",
}

# Chart series colours, in legend order. Cobalt leads, then the three
# harmonized accents, then a light cobalt + neutral slate for "rest".
CHART_PALETTE = [
    RA["signal_600"], RA["gold"], RA["teal"], RA["plum"],
    RA["signal_300"], RA["slate"],
]

# Typography. Brand fonts; HTML preview imports them from Google Fonts and
# PowerPoint substitutes gracefully if they aren't installed locally.
FONT_SANS = "Manrope"
FONT_SERIF = "Playfair Display"   # display: covers, hero numbers, section openers, quotes (swapped from Source Serif 4, 2026-06-11 user request)
FONT_MONO = "JetBrains Mono"

# Skills are bundled via mlflow code_paths (the orchestrator's ``skills`` dir
# is a symlink to ``hackathon-skills``). The deck template — when present —
# resolves from the bundled skill dir.
_SKILL_ROOT = Path(__file__).resolve().parent.parent / "skills" / "compose-pptx"

# Optional RA template. Absent by default → from-scratch mode. Override with
# the ``COMPOSE_DECK_TEMPLATE`` env var or ``template_path=`` on the factory.
_DEFAULT_TEMPLATE_PATH = _SKILL_ROOT / "templates" / "ra_template.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _rgb(hex_: str) -> RGBColor:
    s = hex_.lstrip("#").strip()
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _document_id(title: str, deck_spec: list) -> str:
    canon = json.dumps(deck_spec, sort_keys=True, separators=(",", ":"))
    h = hashlib.sha256(f"deck|{title}|{canon}".encode("utf-8")).hexdigest()[:12]
    return f"deck_{h}"


def _slug(title: str, fallback: str = "deck") -> str:
    s = re.sub(r"[^A-Za-z0-9._-]+", "_", (title or "").strip())
    s = re.sub(r"_+", "_", s).strip("_")
    return (s or fallback)[:80]


# Background semantics. Returns (fill, text, accent) where ``accent`` drives
# the eyebrow dash, bullet markers, timeline bars, and the hero stat colour.
# Light backgrounds accent in cobalt *Signal*; dark backgrounds accent in
# *gold* (amber) — matching the design's `.bg-ink`/`.bg-signal` rules.
#
# Aliases: legacy names ("red", "charcoal", "dark") map onto RA semantics so
# older specs don't crash — "red" → signal (cobalt), "charcoal"/"dark" → ink.
def _bg_colours(bg: str) -> tuple[str, str, str]:
    bg = (bg or "").lower()
    if bg in ("signal", "red"):          # cobalt brand background, white text
        return RA["signal_700"], RA["white"], RA["gold"]
    if bg in ("ink", "charcoal", "dark"):  # near-black cover/closing background
        return RA["ink_950"], RA["ink_100"], RA["gold"]
    if bg == "cream":                    # warm off-white
        return RA["cream"], RA["ink_900"], RA["signal_600"]
    return RA["white"], RA["ink_900"], RA["signal_600"]   # paper (default)


def _is_dark(fill_hex: str) -> bool:
    return fill_hex not in (RA["white"], RA["cream"], RA["paper"])


# ---------------------------------------------------------------------------
# Spec validation
# ---------------------------------------------------------------------------


_VALID_TYPES = {
    "cover", "section_divider", "stat_callout", "kpi_grid", "bullets",
    "two_column", "chart", "chart_commentary", "table", "quote",
    "data_note", "timeline", "closing",
}
_VALID_CHARTS = {"column", "bar", "line", "pie"}


def _coerce_commentary(val) -> list:
    """Normalize the many shapes LLMs emit for a ``chart_commentary`` slide's
    ``commentary`` field into the canonical list-of-``{label, text}`` dicts the
    PPTX + HTML renderers expect.

    The strict schema is a list:
        "commentary": [{"label": "Sweden leads", "text": "..."}, ...]
    But LLMs frequently emit a dict instead — most commonly a
    ``{"title": ..., "bullets": [...]}`` wrapper, or a ``label -> text`` map, or
    a bare string. The renderers slice this (``[:5]``); slicing a dict raises
    ``TypeError: unhashable type: 'slice'`` and crashes the whole render. This
    coercion flattens every variation into a list so the slice + ``.get()``
    loop is safe. Never raises; returns ``[]`` for anything it can't interpret.
    """
    if val is None:
        return []
    if isinstance(val, str):
        return [{"text": val}]
    if isinstance(val, dict):
        for key in ("commentary", "bullets", "points", "items",
                    "paragraphs", "notes", "lines"):
            inner = val.get(key)
            if isinstance(inner, list):
                return _coerce_commentary(inner)
        if "text" in val or "label" in val:
            return [val]
        return [{"label": str(k), "text": v}
                for k, v in val.items() if isinstance(v, str)]
    if isinstance(val, list):
        out = []
        for it in val:
            if isinstance(it, dict):
                out.append(it)
            elif isinstance(it, str):
                out.append({"text": it})
        return out
    return []


def _normalize_slide(sl: dict) -> None:
    """In-place alias normalization. LLMs reuse the kpi_grid key `items` (or
    `points`) for bullet lists; without this the bullets slide renders title-only
    and the content is SILENTLY dropped (the 'last slide is missing information'
    bug). Applies to bullets AND closing (closing next-steps lists)."""
    t = sl.get("type")
    if t in ("bullets", "closing") and not sl.get("bullets"):
        for alias in ("items", "points"):
            v = sl.get(alias)
            if isinstance(v, list) and v and all(isinstance(x, (str, int, float)) for x in v):
                sl["bullets"] = [str(x) for x in v]
                break


def _validate_spec(deck_spec: list) -> tuple[bool, str]:
    if not isinstance(deck_spec, list):
        return False, "deck_spec must be a list of slide dicts"
    if not deck_spec:
        return False, "deck_spec is empty — need at least one slide"

    for i, sl in enumerate(deck_spec):
        if not isinstance(sl, dict):
            return False, f"slide[{i}] is not a dict"
        _normalize_slide(sl)
        t = sl.get("type")
        # isinstance guard first: a dict/list here is unhashable and would
        # crash the `in` membership test (TypeError) instead of returning a
        # correctable validation message to the agent.
        if not isinstance(t, str) or t not in _VALID_TYPES:
            return False, f"slide[{i}].type={t!r} is invalid (must be a string, one of: {sorted(_VALID_TYPES)})"
        if t == "bullets" and not sl.get("bullets"):
            return False, (f"slide[{i}] type 'bullets' has no bullet content — put the "
                           f"list of strings in `bullets` (aliases `items`/`points` are accepted)")
        if t in ("chart", "chart_commentary"):
            ct = sl.get("chart")
            if not isinstance(ct, str) or ct not in _VALID_CHARTS:
                return False, (
                    f"slide[{i}].chart={ct!r} is invalid — `chart` must be a STRING "
                    f"chart-type, one of: {sorted(_VALID_CHARTS)} (chart data goes in "
                    f"the slide's own `categories` + `series` keys, not inside `chart`)"
                )
            cats = sl.get("categories") or []
            series = sl.get("series") or []
            if not isinstance(cats, list) or not isinstance(series, list) or not cats or not series:
                return False, f"slide[{i}] chart needs non-empty `categories` (list) + `series` (list)"
            for j, s in enumerate(series):
                vals = s.get("values") if isinstance(s, dict) else None
                if not isinstance(vals, list) or len(vals) != len(cats):
                    return False, (
                        f"slide[{i}].series[{j}].values must be a list matching categories "
                        f"(got {len(vals) if isinstance(vals, list) else vals!r} vs {len(cats)})"
                    )

    return True, ""


# ---------------------------------------------------------------------------
# PPTX rendering — low-level shape helpers
# ---------------------------------------------------------------------------


def _set_fill(shape, hex_: str) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_)


def _no_line(shape) -> None:
    shape.line.fill.background()


def _add_textbox(slide, x_in, y_in, w_in, h_in, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return tb, tf


def _run(p, text: str, *, size_pt: float, color_hex: str, bold: bool = False,
         italic: bool = False, font: str = FONT_SANS, tracking: float | None = None,
         align: PP_ALIGN | None = None):
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size_pt)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = _rgb(color_hex)
    return r


# Template mode: when True, helpers that paint full-bleed backgrounds, the
# brand-mark, the cover orb, and the page-number footer all early-return — the
# loaded ``ra_template.pptx`` already supplies that chrome via its master +
# layouts.
_TEMPLATE_MODE: contextvars.ContextVar[bool] = contextvars.ContextVar(
    "compose_deck_template_mode", default=False,
)


def _add_background(slide, hex_: str) -> None:
    if _TEMPLATE_MODE.get():
        return
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
    )
    _set_fill(rect, hex_)
    _no_line(rect)
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)


def _draw_eyebrow(slide, text: str, x_in: float, y_in: float, color_hex: str) -> None:
    """Eyebrow = a short dash + uppercase tracked label. The dash + label are
    the RA system's single accent line — never an underline below the title."""
    if not text:
        return
    dash = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x_in), Inches(y_in + 0.06),
        Inches(0.25), Inches(0.025),
    )
    _set_fill(dash, color_hex)
    _no_line(dash)
    tb, tf = _add_textbox(slide, x_in + 0.32, y_in - 0.02, 9.0, 0.3)
    p = tf.paragraphs[0]
    _run(p, text.upper(), size_pt=10, color_hex=color_hex, bold=True, tracking=0.16)


# ---------------------------------------------------------------------------
# RA brand-mark + cover orb (drawn natively — no raster logos)
# ---------------------------------------------------------------------------


def _draw_brandmark(slide, x_in: float, y_in: float, *, on_dark: bool,
                    wordmark: bool = True, mk_in: float = 0.30) -> None:
    """Draw the Resonance Analytics brand-mark: a rounded *Signal* square with
    a small cross-hair, optionally followed by the "Resonance Analytics"
    wordmark. On dark backgrounds the square inverts to white. Mirrors the
    design system's ``.brand-mark .mk`` glyph."""
    if _TEMPLATE_MODE.get():
        return
    sq_hex = RA["white"] if on_dark else RA["signal_600"]
    notch_hex = RA["signal_700"] if on_dark else RA["white"]
    text_hex = RA["white"] if on_dark else RA["ink_900"]

    sq = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_in), Inches(y_in), Inches(mk_in), Inches(mk_in),
    )
    try:
        sq.adjustments[0] = 0.22
    except Exception:  # noqa: BLE001
        pass
    _set_fill(sq, sq_hex)
    _no_line(sq)

    # Cross-hair: a short horizontal bar (upper-right) + vertical bar
    # (upper-left) — the two notches of the glyph.
    bar = mk_in * 0.30
    thick = mk_in * 0.075
    h = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x_in + mk_in * 0.40), Inches(y_in + mk_in * 0.28),
        Inches(bar), Inches(thick),
    )
    _set_fill(h, notch_hex)
    _no_line(h)
    v = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x_in + mk_in * 0.28), Inches(y_in + mk_in * 0.28),
        Inches(thick), Inches(bar * 1.25),
    )
    _set_fill(v, notch_hex)
    _no_line(v)

    if wordmark:
        tb, tf = _add_textbox(slide, x_in + mk_in + 0.14, y_in - 0.05,
                              5.0, mk_in + 0.1, anchor=MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], "Resonance Analytics", size_pt=14,
             color_hex=text_hex, bold=True)


def _grad_oval(slide, x_in, y_in, d_in, hex_from, hex_to, angle=90.0):
    """An oval with a two-stop linear gradient — used to fake the soft radial
    glows of the cover poster (python-pptx has no radial gradient). Fading a
    colour into the surrounding background colour reads as a glow rather than
    a flat disc."""
    ov = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x_in), Inches(y_in), Inches(d_in), Inches(d_in),
    )
    try:
        g = ov.fill
        g.gradient()
        stops = g.gradient_stops
        stops[0].color.rgb = _rgb(hex_from)
        stops[1].color.rgb = _rgb(hex_to)
        try:
            g.gradient_angle = angle
        except Exception:  # noqa: BLE001
            pass
    except Exception:  # noqa: BLE001
        _set_fill(ov, hex_from)
    _no_line(ov)
    return ov


def _draw_cover_orb(slide) -> None:
    """The cover's signature cobalt *orb* — a glowing oval in the upper-right,
    with an amber rim highlight and a soft amber under-glow bottom-left. Built
    from layered gradient ovals: the orb fades Signal-500 → Signal-900, the
    rim highlight fades amber → Signal-800 (so it blends as a highlight), and
    the corner under-glow fades amber → the near-black cover base (so it reads
    as a glow, not a flat disc)."""
    if _TEMPLATE_MODE.get():
        return
    # Large cobalt orb, top-right, partly off-canvas.
    orb_d = 6.4
    ox = SLIDE_W_IN - orb_d + 1.7
    oy = -1.7
    _grad_oval(slide, ox, oy, orb_d, RA["signal_500"], RA["signal_900"], angle=135.0)

    # Amber rim highlight on the upper-left of the orb — fades into the orb.
    hi_d = 2.4
    _grad_oval(slide, ox + orb_d * 0.10, oy + orb_d * 0.10, hi_d,
               RA["amber_light"], RA["signal_800"], angle=135.0)

    # Soft amber under-glow, bottom-left — fades into the near-black base.
    gl_d = 3.4
    _grad_oval(slide, -1.9, SLIDE_H_IN - gl_d + 2.0, gl_d,
               RA["gold_deep"], RA["cover_base"], angle=315.0)


def _draw_slide_foot(slide, n_str: str, footer_note: str, text_hex: str, num_hex: str,
                     *, brandmark_on_dark: bool | None = None) -> None:
    """Footer at bottom: optional left note + page number on the right, and
    (on content slides) a tiny RA glyph to the left of the page number."""
    if _TEMPLATE_MODE.get():
        return
    tb_l, tf_l = _add_textbox(slide, 0.75, SLIDE_H_IN - 0.46, 7.0, 0.3)
    _run(tf_l.paragraphs[0], footer_note or "", size_pt=10, color_hex=text_hex)
    tb_r, tf_r = _add_textbox(slide, SLIDE_W_IN - 2.0, SLIDE_H_IN - 0.46, 1.25, 0.3)
    _run(tf_r.paragraphs[0], n_str, size_pt=10, color_hex=num_hex, bold=True,
         font=FONT_MONO, align=PP_ALIGN.RIGHT)
    if brandmark_on_dark is not None:
        # Tiny mark (no wordmark) just left of the page number.
        _draw_brandmark(slide, SLIDE_W_IN - 2.0 - 0.30 - 0.16,
                        SLIDE_H_IN - 0.49, on_dark=brandmark_on_dark,
                        wordmark=False, mk_in=0.26)


# ---------------------------------------------------------------------------
# PPTX rendering — slide renderers
# ---------------------------------------------------------------------------


def _render_cover(slide, sl: dict, idx: int, total: int, fill_hex: str,
                  text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    # Cobalt orb behind the title (added FIRST so text overlays it). Skipped
    # on light covers or if the agent opts out via `motif: false`.
    if on_dark and sl.get("motif", True):
        _draw_cover_orb(slide)
    # Brand-mark + wordmark, top-left.
    _draw_brandmark(slide, 0.75, 0.62, on_dark=on_dark)
    _draw_eyebrow(slide, sl.get("eyebrow", ""), 0.75, 2.7, accent_hex)
    tb, tf = _add_textbox(slide, 0.75, 3.0, SLIDE_W_IN - 1.5, 2.6)
    p = tf.paragraphs[0]
    # Display serif for the cover hero (the design's signature).
    _run(p, sl.get("title", ""), size_pt=72, color_hex=text_hex, font=FONT_SERIF)
    if sl.get("lede"):
        muted = RA["ink_300"] if on_dark else RA["ink_600"]
        tb2, tf2 = _add_textbox(slide, 0.75, 5.5, SLIDE_W_IN - 1.5, 1.0)
        _run(tf2.paragraphs[0], sl["lede"], size_pt=22, color_hex=muted)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     text_hex, text_hex)


def _render_section_divider(slide, sl: dict, idx: int, total: int,
                            fill_hex: str, text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    _draw_eyebrow(slide, sl.get("eyebrow", ""), 0.75, 3.0, accent_hex)
    tb, tf = _add_textbox(slide, 0.75, 3.3, SLIDE_W_IN - 1.5, 1.6)
    _run(tf.paragraphs[0], sl.get("title", ""), size_pt=56, color_hex=text_hex,
         font=FONT_SERIF)
    if sl.get("lede"):
        muted = RA["ink_300"] if on_dark else RA["ink_600"]
        tb2, tf2 = _add_textbox(slide, 0.75, 4.8, SLIDE_W_IN - 1.5, 1.0)
        _run(tf2.paragraphs[0], sl["lede"], size_pt=22, color_hex=muted)
    _draw_brandmark(slide, 0.75, 0.62, on_dark=on_dark, wordmark=False)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     text_hex, text_hex)


def _slide_head(slide, sl: dict, text_hex: str, accent_hex: str) -> float:
    """Eyebrow + sans title (+ optional lede). Returns the y where body starts.
    Content-slide titles are Manrope (sans) 600 — serif is reserved for
    display moments (cover, section, hero stat, KPI value, quote)."""
    _draw_eyebrow(slide, sl.get("eyebrow", ""), 0.75, 0.65, accent_hex)
    tb, tf = _add_textbox(slide, 0.75, 0.92, SLIDE_W_IN - 1.5, 0.9)
    _run(tf.paragraphs[0], sl.get("title", ""), size_pt=40, color_hex=text_hex, bold=True)
    y = 1.85
    if sl.get("lede"):
        tb2, tf2 = _add_textbox(slide, 0.75, y, SLIDE_W_IN - 1.5, 0.6)
        _run(tf2.paragraphs[0], sl["lede"], size_pt=18, color_hex=RA["ink_500"])
        y += 0.6
    return y + 0.2


def _render_stat_callout(slide, sl: dict, idx: int, total: int,
                         fill_hex: str, text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    y = _slide_head(slide, sl, text_hex, accent_hex)
    value = str(sl.get("value", ""))
    unit = sl.get("unit", "")
    # One emphatic brand moment per deck — the hero number in cobalt (or gold
    # on dark), set in the display serif.
    val_hex = accent_hex
    unit_hex = RA["ink_400"] if on_dark else RA["ink_600"]
    tb, tf = _add_textbox(slide, 0.75, y + 0.4, SLIDE_W_IN - 1.5, 2.2)
    p = tf.paragraphs[0]
    _run(p, value, size_pt=130, color_hex=val_hex, font=FONT_SERIF)
    if unit:
        _run(p, " " + unit, size_pt=44, color_hex=unit_hex, bold=True)
    if sl.get("delta"):
        delta_dir = sl.get("delta_dir", "flat")
        delta_color = (RA["success"] if delta_dir == "up" else
                       RA["danger"]  if delta_dir == "down" else
                       RA["ink_500"])
        arrow = "▲ " if delta_dir == "up" else ("▼ " if delta_dir == "down" else "")
        tb_d, tf_d = _add_textbox(slide, 0.75, y + 2.9, 6.0, 0.4)
        _run(tf_d.paragraphs[0], arrow + sl["delta"], size_pt=22,
             color_hex=delta_color, bold=True)
    if sl.get("caption"):
        cap_hex = RA["ink_300"] if on_dark else RA["ink_500"]
        tb_c, tf_c = _add_textbox(slide, 0.75, SLIDE_H_IN - 1.4, SLIDE_W_IN - 1.5, 0.7)
        _run(tf_c.paragraphs[0], sl["caption"], size_pt=15, color_hex=cap_hex)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     RA["ink_500"] if not on_dark else RA["ink_300"],
                     text_hex, brandmark_on_dark=on_dark)


def _render_kpi_grid(slide, sl: dict, idx: int, total: int,
                     fill_hex: str, text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    y = _slide_head(slide, sl, text_hex, accent_hex)
    kpis = sl.get("kpis", [])[:6]
    n = len(kpis)
    if n == 0:
        return
    if n <= 4:
        cols, rows = n, 1
    else:
        cols, rows = 3, 2
    gutter = 0.3
    avail_w = SLIDE_W_IN - 1.5
    card_w = (avail_w - gutter * (cols - 1)) / cols
    card_h = 2.4 if rows == 1 else 1.9
    start_y = y + 0.2

    card_fill = RA["ink_900"] if on_dark else RA["white"]
    border_hex = RA["ink_700"] if on_dark else RA["ink_200"]
    label_hex = RA["ink_400"] if on_dark else RA["ink_600"]
    val_hex = RA["white"] if on_dark else RA["ink_900"]

    for i, k in enumerate(kpis):
        r = i // cols
        c = i % cols
        x = 0.75 + c * (card_w + gutter)
        cy = start_y + r * (card_h + gutter)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x), Inches(cy),
                                      Inches(card_w), Inches(card_h))
        try:
            card.adjustments[0] = 0.06
        except Exception:  # noqa: BLE001
            pass
        _set_fill(card, card_fill)
        card.line.color.rgb = _rgb(border_hex)
        card.line.width = Pt(1.0)
        # Accent tick at the top-left of each tile.
        tick = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(x + 0.25), Inches(cy + 0.25),
                                      Inches(0.4), Inches(0.04))
        _set_fill(tick, accent_hex)
        _no_line(tick)
        tb_l, tf_l = _add_textbox(slide, x + 0.25, cy + 0.42, card_w - 0.5, 0.35)
        _run(tf_l.paragraphs[0], (k.get("label", "") or "").upper(), size_pt=11,
             color_hex=label_hex, bold=True)
        val = str(k.get("value", ""))
        unit = k.get("unit", "")
        tb_v, tf_v = _add_textbox(slide, x + 0.25, cy + 0.86, card_w - 0.5, 1.1)
        p = tf_v.paragraphs[0]
        # KPI values in the display serif (supporting numbers, ink-coloured).
        _run(p, val, size_pt=44, color_hex=val_hex, font=FONT_SERIF)
        if unit:
            _run(p, " " + unit, size_pt=20, color_hex=label_hex, bold=True)
        if k.get("delta"):
            dd = k.get("delta_dir", "flat")
            dc = (RA["success"] if dd == "up" else
                  RA["danger"]  if dd == "down" else label_hex)
            arrow = "▲ " if dd == "up" else ("▼ " if dd == "down" else "")
            tb_d, tf_d = _add_textbox(slide, x + 0.25, cy + card_h - 0.55,
                                      card_w - 0.5, 0.35)
            _run(tf_d.paragraphs[0], arrow + k["delta"], size_pt=13,
                 color_hex=dc, bold=True)
    if sl.get("caption"):
        cap_hex = RA["ink_300"] if on_dark else RA["ink_500"]
        tb_c, tf_c = _add_textbox(slide, 0.75, SLIDE_H_IN - 1.1, SLIDE_W_IN - 1.5, 0.5)
        _run(tf_c.paragraphs[0], sl["caption"], size_pt=14, color_hex=cap_hex)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     RA["ink_500"] if not on_dark else RA["ink_300"],
                     text_hex, brandmark_on_dark=on_dark)


def _render_bullets(slide, sl: dict, idx: int, total: int,
                    fill_hex: str, text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    y = _slide_head(slide, sl, text_hex, accent_hex)
    bullets = sl.get("bullets", [])
    tb, tf = _add_textbox(slide, 0.75, y + 0.3, SLIDE_W_IN - 1.5, SLIDE_H_IN - y - 1.5)
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _run(p, "•  ", size_pt=22, color_hex=accent_hex, bold=True)
        _run(p, str(b), size_pt=22, color_hex=text_hex)
        p.space_after = Pt(14)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     RA["ink_500"] if not on_dark else RA["ink_300"],
                     text_hex, brandmark_on_dark=on_dark)


def _render_two_column(slide, sl: dict, idx: int, total: int,
                       fill_hex: str, text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    y = _slide_head(slide, sl, text_hex, accent_hex)
    avail_w = SLIDE_W_IN - 1.5
    col_w = (avail_w - 0.5) / 2
    for i, side in enumerate(("left", "right")):
        block = sl.get(side, {}) or {}
        x = 0.75 + i * (col_w + 0.5)
        tb_h, tf_h = _add_textbox(slide, x, y + 0.2, col_w, 0.5)
        _run(tf_h.paragraphs[0], block.get("heading", ""), size_pt=24,
             color_hex=text_hex, bold=True)
        bullets = block.get("bullets", [])
        tb_b, tf_b = _add_textbox(slide, x, y + 0.85, col_w, SLIDE_H_IN - y - 2.1)
        first = True
        for b in bullets:
            p = tf_b.paragraphs[0] if first else tf_b.add_paragraph()
            first = False
            _run(p, "•  ", size_pt=18, color_hex=accent_hex, bold=True)
            _run(p, str(b), size_pt=18, color_hex=text_hex)
            p.space_after = Pt(10)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     RA["ink_500"] if not on_dark else RA["ink_300"],
                     text_hex, brandmark_on_dark=on_dark)


_PPTX_CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
    "line":   XL_CHART_TYPE.LINE,
    "pie":    XL_CHART_TYPE.PIE,
}


def _add_native_chart(slide, sl: dict, x_in: float, y_in: float,
                      w_in: float, h_in: float, show_legend: bool,
                      on_dark: bool = False) -> None:
    chart_kind = sl.get("chart", "column")
    label_hex = RA["white"] if on_dark else RA["ink_900"]
    chart_data = CategoryChartData()
    cats = [str(c) for c in (sl.get("categories") or [])]
    chart_data.categories = cats
    series = sl.get("series", [])
    for s in series:
        chart_data.add_series(s.get("name", ""), list(s.get("values", [])))

    xl_kind = _PPTX_CHART_TYPES[chart_kind]
    chart_shape = slide.shapes.add_chart(
        xl_kind, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_title = False
    if show_legend and len(series) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        try:
            chart.legend.font.size = Pt(11)
            chart.legend.font.color.rgb = _rgb(label_hex)
            chart.legend.font.name = FONT_SANS
        except Exception:  # noqa: BLE001
            pass
    else:
        chart.has_legend = False

    # On dark backgrounds, force white tick labels so they don't vanish.
    # NOTE: pie charts have no category/value axis — python-pptx raises
    # ValueError (not AttributeError) on access, so getattr(default) does NOT
    # shield it. This crashed dark-bg pie slides ("chart has no category axis").
    if on_dark:
        for ax_name in ("category_axis", "value_axis"):
            try:
                ax = getattr(chart, ax_name)
                ax.tick_labels.font.color.rgb = _rgb(RA["white"])
                ax.tick_labels.font.size = Pt(11)
                ax.tick_labels.font.name = FONT_SANS
            except Exception:  # noqa: BLE001 — no such axis for this chart kind
                pass

    # value_label → value-axis title (column/bar/line have a value axis; pie
    # doesn't). Omitted in the HTML preview, which is approximate.
    value_label = sl.get("value_label")
    if value_label and chart_kind != "pie":
        try:
            va = chart.value_axis
            va.has_title = True
            va.axis_title.text_frame.text = str(value_label)
            for para in va.axis_title.text_frame.paragraphs:
                for r in para.runs:
                    r.font.size = Pt(11)
                    r.font.name = FONT_SANS
                    r.font.color.rgb = _rgb(label_hex)
        except Exception:  # noqa: BLE001
            pass

    plots = list(chart.plots)
    if plots:
        plot = plots[0]
        if chart_kind in ("column", "bar"):
            plot.gap_width = 60

        if chart_kind == "pie" and plot.series:
            ser = list(plot.series)[0]
            for i, pt in enumerate(ser.points):
                color = CHART_PALETTE[i % len(CHART_PALETTE)]
                try:
                    fill = pt.format.fill
                    fill.solid()
                    fill.fore_color.rgb = _rgb(color)
                    pt.format.line.fill.background()
                except Exception:  # noqa: BLE001
                    pass
            try:
                dl = ser.data_labels
                dl.show_percentage = True
                dl.show_category_name = True
                dl.font.size = Pt(11)
                dl.font.name = FONT_SANS
                dl.font.color.rgb = _rgb(label_hex)
            except Exception:  # noqa: BLE001
                pass
        else:
            for i, ser in enumerate(plot.series):
                color = CHART_PALETTE[i % len(CHART_PALETTE)]
                try:
                    if chart_kind == "line":
                        line = ser.format.line
                        line.color.rgb = _rgb(color)
                        line.width = Pt(2.5)
                        try:
                            mkr = ser.marker
                            mkr.style = 8
                            mkr.size = 7
                            mfill = mkr.format.fill
                            mfill.solid()
                            mfill.fore_color.rgb = _rgb(color)
                            mkr.format.line.color.rgb = _rgb(color)
                        except Exception:  # noqa: BLE001
                            pass
                    else:
                        fill = ser.format.fill
                        fill.solid()
                        fill.fore_color.rgb = _rgb(color)
                        ser.format.line.fill.background()
                except Exception:  # noqa: BLE001
                    pass
                if sl.get("show_values", True) and chart_kind != "line":
                    try:
                        dl = ser.data_labels
                        dl.show_value = True
                        if chart_kind in ("column", "bar"):
                            dl.position = XL_LABEL_POSITION.OUTSIDE_END
                        dl.font.size = Pt(10)
                        dl.font.name = FONT_SANS
                        dl.font.color.rgb = _rgb(label_hex)
                    except Exception:  # noqa: BLE001
                        pass


def _render_chart(slide, sl: dict, idx: int, total: int,
                  fill_hex: str, text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    cap_hex = RA["ink_300"] if on_dark else RA["ink_500"]
    foot_note_hex = RA["ink_400"] if on_dark else RA["ink_500"]
    y = _slide_head(slide, sl, text_hex, accent_hex)
    _add_native_chart(slide, sl,
                      x_in=0.75, y_in=y + 0.2,
                      w_in=SLIDE_W_IN - 1.5, h_in=SLIDE_H_IN - y - 1.7,
                      show_legend=bool(sl.get("show_legend", False)),
                      on_dark=on_dark)
    if sl.get("caption"):
        tb_c, tf_c = _add_textbox(slide, 0.75, SLIDE_H_IN - 1.0, SLIDE_W_IN - 1.5, 0.5)
        _run(tf_c.paragraphs[0], sl["caption"], size_pt=14, color_hex=cap_hex)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     foot_note_hex, text_hex, brandmark_on_dark=on_dark)


def _render_chart_commentary(slide, sl: dict, idx: int, total: int,
                             fill_hex: str, text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    y = _slide_head(slide, sl, text_hex, accent_hex)
    chart_w = (SLIDE_W_IN - 1.5) * 0.60
    side_w = (SLIDE_W_IN - 1.5) * 0.36
    _add_native_chart(slide, sl, x_in=0.75, y_in=y + 0.2,
                      w_in=chart_w, h_in=SLIDE_H_IN - y - 1.7,
                      show_legend=bool(sl.get("show_legend", False)),
                      on_dark=on_dark)
    side_x = 0.75 + chart_w + 0.3
    body_hex = RA["ink_200"] if on_dark else RA["ink_700"]
    tb, tf = _add_textbox(slide, side_x, y + 0.2, side_w, SLIDE_H_IN - y - 1.5)
    items = _coerce_commentary(sl.get("commentary"))[:5]
    first = True
    for item in items:
        if not isinstance(item, dict):
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _run(p, item.get("label", ""), size_pt=14, color_hex=accent_hex, bold=True)
        _run(p, "  " + item.get("text", ""), size_pt=14, color_hex=body_hex)
        p.space_after = Pt(12)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     RA["ink_500"] if not on_dark else RA["ink_300"],
                     text_hex, brandmark_on_dark=on_dark)


def _render_table(slide, sl: dict, idx: int, total: int,
                  fill_hex: str, text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    y = _slide_head(slide, sl, text_hex, accent_hex)
    headers = sl.get("headers", [])
    rows = sl.get("rows", [])
    if not headers or not rows:
        return
    n_rows = len(rows) + 1
    n_cols = len(headers)
    avail_h = SLIDE_H_IN - y - 1.6
    row_h = min(0.6, avail_h / n_rows)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                       Inches(0.75), Inches(y + 0.2),
                                       Inches(SLIDE_W_IN - 1.5),
                                       Inches(row_h * n_rows))
    table = tbl_shape.table
    # Header row in the cobalt Signal primary.
    band_a = RA["white"] if not on_dark else RA["ink_900"]
    band_b = RA["ink_100"] if not on_dark else RA["ink_800"]
    body_text = RA["ink_900"] if not on_dark else RA["ink_100"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(RA["signal_600"])
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        _run(p, str(h), size_pt=14, color_hex=RA["white"], bold=True)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            if c >= n_cols:
                break
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(band_a if r % 2 == 0 else band_b)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            s = str(val)
            header = str(headers[c])
            is_num_col = any(t in header.lower() for t in ("%", "thb", "usd",
                                                            "eur", "gbp",
                                                            "(m)", "delta", "δ"))
            if is_num_col:
                p.alignment = PP_ALIGN.RIGHT
            _run(p, s, size_pt=13, color_hex=body_text)
    if sl.get("caption"):
        cap_hex = RA["ink_300"] if on_dark else RA["ink_500"]
        tb_c, tf_c = _add_textbox(slide, 0.75, SLIDE_H_IN - 1.0, SLIDE_W_IN - 1.5, 0.5)
        _run(tf_c.paragraphs[0], sl["caption"], size_pt=13, color_hex=cap_hex)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     RA["ink_500"] if not on_dark else RA["ink_300"],
                     text_hex, brandmark_on_dark=on_dark)


def _render_quote(slide, sl: dict, idx: int, total: int,
                  fill_hex: str, text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    _draw_eyebrow(slide, sl.get("eyebrow", ""), 0.75, 1.2, accent_hex)
    tb, tf = _add_textbox(slide, 1.5, 1.8, SLIDE_W_IN - 3.0, 3.6,
                          anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    # Pull-quote in editorial serif italic.
    _run(p, "“" + str(sl.get("quote", "")) + "”", size_pt=40,
         color_hex=text_hex, italic=True, font=FONT_SERIF)
    muted = RA["ink_400"] if on_dark else RA["ink_600"]
    if sl.get("attribution"):
        tb_a, tf_a = _add_textbox(slide, 1.5, 5.7, SLIDE_W_IN - 3.0, 0.5)
        _run(tf_a.paragraphs[0], "— " + sl["attribution"], size_pt=18,
             color_hex=muted, bold=True)
    if sl.get("caption"):
        tb_c, tf_c = _add_textbox(slide, 1.5, 6.3, SLIDE_W_IN - 3.0, 0.5)
        _run(tf_c.paragraphs[0], sl["caption"], size_pt=14, color_hex=muted)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     RA["ink_500"] if not on_dark else RA["ink_300"],
                     text_hex, brandmark_on_dark=on_dark)


def _render_data_note(slide, sl: dict, idx: int, total: int,
                      fill_hex: str, text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    y = _slide_head(slide, sl, text_hex, accent_hex)
    body = (
        f"Requested: {sl.get('requested', '—')}.  "
        f"Available data: {sl.get('available', '—')}.  "
        f"{sl.get('reason', '')}"
    )
    body_hex = RA["ink_200"] if on_dark else RA["ink_700"]
    tb, tf = _add_textbox(slide, 0.75, y + 0.3, SLIDE_W_IN - 1.5, 3.0)
    _run(tf.paragraphs[0], body, size_pt=22, color_hex=body_hex)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     RA["ink_500"] if not on_dark else RA["ink_300"],
                     text_hex, brandmark_on_dark=on_dark)


def _timeline_heading_pt(text: str, col_w_in: float) -> float:
    """Shrink the timeline-step heading font when the text would overflow the
    column. Bands keep ≤2 lines inside the heading box (1.2" tall)."""
    n = len(text or "")
    cap_22 = int(col_w_in * 11)
    cap_18 = int(col_w_in * 14)
    cap_16 = int(col_w_in * 16)
    if n <= cap_22:
        return 22
    if n <= cap_18:
        return 18
    if n <= cap_16:
        return 16
    return 14


def _render_timeline(slide, sl: dict, idx: int, total: int,
                     fill_hex: str, text_hex: str, accent_hex: str) -> None:
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    body_hex = RA["ink_200"] if on_dark else RA["ink_700"]
    y = _slide_head(slide, sl, text_hex, accent_hex)
    steps = sl.get("steps", [])[:5]
    n = len(steps)
    if not n:
        return
    avail_w = SLIDE_W_IN - 1.5
    gutter = 0.3
    col_w = (avail_w - gutter * (n - 1)) / n
    head_h = 1.2
    body_y_offset = 2.05
    col_h = SLIDE_H_IN - y - 1.5
    for i, step in enumerate(steps):
        x = 0.75 + i * (col_w + gutter)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(x), Inches(y + 0.2),
                                     Inches(col_w), Inches(0.05))
        _set_fill(bar, accent_hex)
        _no_line(bar)
        tb_l, tf_l = _add_textbox(slide, x, y + 0.35, col_w, 0.35)
        _run(tf_l.paragraphs[0], (step.get("label", "") or "").upper(), size_pt=11,
             color_hex=accent_hex, bold=True)
        heading = step.get("heading", "")
        h_pt = _timeline_heading_pt(heading, col_w)
        tb_h, tf_h = _add_textbox(slide, x, y + 0.75, col_w, head_h)
        _run(tf_h.paragraphs[0], heading, size_pt=h_pt, color_hex=text_hex, bold=True)
        tb_b, tf_b = _add_textbox(slide, x, y + body_y_offset, col_w,
                                  col_h - body_y_offset + 0.4)
        _run(tf_b.paragraphs[0], step.get("body", ""), size_pt=14, color_hex=body_hex)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     RA["ink_500"] if not on_dark else RA["ink_300"],
                     text_hex, brandmark_on_dark=on_dark)


def _render_closing(slide, sl: dict, idx: int, total: int,
                    fill_hex: str, text_hex: str, accent_hex: str) -> None:
    # Closing mirrors the cover in inverse — default to the dark ink poster.
    if sl.get("bg") is None:
        fill_hex, text_hex, accent_hex = _bg_colours("ink")
    _add_background(slide, fill_hex)
    on_dark = _is_dark(fill_hex)
    if on_dark and sl.get("motif", True):
        _draw_cover_orb(slide)
    _draw_brandmark(slide, 0.75, 0.62, on_dark=on_dark)
    _draw_eyebrow(slide, sl.get("eyebrow", "") or "NEXT STEPS", 0.75, 2.7, accent_hex)
    tb, tf = _add_textbox(slide, 0.75, 3.0, SLIDE_W_IN - 1.5, 2.0)
    _run(tf.paragraphs[0], sl.get("title", ""), size_pt=56, color_hex=text_hex,
         font=FONT_SERIF)
    if sl.get("lede"):
        muted = RA["ink_300"] if on_dark else RA["ink_600"]
        tb2, tf2 = _add_textbox(slide, 0.75, 4.9, SLIDE_W_IN - 1.5, 1.0)
        _run(tf2.paragraphs[0], sl["lede"], size_pt=22, color_hex=muted)
    if sl.get("bullets"):
        # closing next-steps list (normalized from items/points) — never drop content
        muted = RA["ink_300"] if on_dark else RA["ink_600"]
        y0 = 5.6 if sl.get("lede") else 4.9
        tb3, tf3 = _add_textbox(slide, 0.75, y0, SLIDE_W_IN - 1.5, SLIDE_H_IN - y0 - 0.7)
        first = True
        for b in sl["bullets"]:
            p = tf3.paragraphs[0] if first else tf3.add_paragraph()
            first = False
            _run(p, "•  ", size_pt=16, color_hex=accent_hex, bold=True)
            _run(p, str(b), size_pt=16, color_hex=muted)
            p.space_after = Pt(8)
    _draw_slide_foot(slide, f"{idx:02d} / {total:02d}", sl.get("footer_note", ""),
                     text_hex, text_hex)


_RENDERERS = {
    "cover":            _render_cover,
    "section_divider":  _render_section_divider,
    "stat_callout":     _render_stat_callout,
    "kpi_grid":         _render_kpi_grid,
    "bullets":          _render_bullets,
    "two_column":       _render_two_column,
    "chart":            _render_chart,
    "chart_commentary": _render_chart_commentary,
    "table":            _render_table,
    "quote":            _render_quote,
    "data_note":        _render_data_note,
    "timeline":         _render_timeline,
    "closing":          _render_closing,
}


# ---------------------------------------------------------------------------
# Template mode (optional — when ra_template.pptx is supplied)
# ---------------------------------------------------------------------------
#
# These are the layout names the renderer looks for inside ra_template.pptx.
# When you build the template in PowerPoint, name your slide layouts to match
# (Slide Master view → rename layouts). See templates/README.md.
_TEMPLATE_COVER_LAYOUT = "RA_TitleSlide"
_TEMPLATE_CONTENT_LAYOUT = "RA_blank"
_TEMPLATE_THANKYOU_LAYOUT = "RA_thankyou"
_TEMPLATE_DARK_CONTENT_LAYOUT = "RA_dark"
_TEMPLATE_SECTION_SIGNAL_LAYOUT = "RA_SectionSignal"
_TEMPLATE_SECTION_WHITE_LAYOUT = "RA_SectionWhite"
_TEMPLATE_CREAM_LAYOUT = "RA_cream"


def _find_layout(pres, name: str):
    for lay in pres.slide_layouts:
        if lay.name == name:
            return lay
    return None


def _pick_template_layout(pres, sl_type: str, bg: str, default_content_layout):
    # Cream always wins (content + section_divider) so bg:"cream" stops
    # silently rendering on the white RA_blank layout. Falls back to the
    # default layout if RA_cream is absent (older templates).
    if bg == "cream":
        lay = _find_layout(pres, _TEMPLATE_CREAM_LAYOUT)
        if lay is not None:
            return lay
    if sl_type == "section_divider":
        if bg in ("signal", "red"):
            lay = _find_layout(pres, _TEMPLATE_SECTION_SIGNAL_LAYOUT)
            if lay is not None:
                return lay
        if bg in ("white", "paper"):
            lay = _find_layout(pres, _TEMPLATE_SECTION_WHITE_LAYOUT)
            if lay is not None:
                return lay
    if bg in ("ink", "charcoal", "dark"):
        lay = _find_layout(pres, _TEMPLATE_DARK_CONTENT_LAYOUT)
        if lay is not None:
            return lay
    return default_content_layout


def _set_placeholder_text(slide, *, ph_idx: int, text: str, size_pt: float,
                          color_hex: str, bold: bool = False,
                          font: str = FONT_SANS) -> None:
    if not text:
        return
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx != ph_idx:
                continue
        except Exception:  # noqa: BLE001
            continue
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        _run(p, text, size_pt=size_pt, color_hex=color_hex, bold=bold, font=font)
        return


def _strip_empty_placeholders(slide) -> None:
    from pptx.enum.shapes import PP_PLACEHOLDER
    sp_tree = slide.shapes._spTree
    to_remove = []
    for ph in list(slide.placeholders):
        try:
            ptype = ph.placeholder_format.type
        except Exception:  # noqa: BLE001
            continue
        if ptype == PP_PLACEHOLDER.TITLE:
            continue
        if not (ph.text_frame.text or "").strip():
            to_remove.append(ph)
    for ph in to_remove:
        sp_tree.remove(ph._element)


def _remove_slide(pres, slide_index: int) -> None:
    sldIdLst = pres.slides._sldIdLst
    children = list(sldIdLst)
    if slide_index < 0 or slide_index >= len(children):
        return
    rId = children[slide_index].rId
    pres.part.drop_rel(rId)
    sldIdLst.remove(children[slide_index])


def _autosize_title_pt(text: str, *, max_pt: float, min_pt: float = 24.0) -> float:
    n = len(text or "")
    if n <= 30:
        return max_pt
    if n <= 50:
        return max(min_pt, max_pt * 0.78)
    if n <= 80:
        return max(min_pt, max_pt * 0.62)
    return min_pt


def _render_template_cover(slide, sl: dict) -> None:
    title = sl.get("title", "")
    _set_placeholder_text(slide, ph_idx=10, text=title,
                          size_pt=_autosize_title_pt(title, max_pt=44),
                          color_hex=RA["white"], bold=False, font=FONT_SERIF)
    subtitle = sl.get("lede") or sl.get("eyebrow") or sl.get("subtitle") or ""
    _set_placeholder_text(slide, ph_idx=1, text=subtitle,
                          size_pt=20, color_hex=RA["white"])


def _render_template_thankyou(slide, sl: dict) -> None:
    title = sl.get("title", "") or "Thank you"
    _set_placeholder_text(slide, ph_idx=11, text=title,
                          size_pt=_autosize_title_pt(title, max_pt=56),
                          color_hex=RA["white"], bold=False, font=FONT_SERIF)
    subtitle = sl.get("lede") or sl.get("eyebrow") or sl.get("footer_note") or ""
    _set_placeholder_text(slide, ph_idx=1, text=subtitle,
                          size_pt=20, color_hex=RA["white"])


def _render_pptx(title: str, deck_spec: list,
                 template_path: str | None = None) -> bytes:
    """Render the deck spec to PPTX bytes.

    ``template_path`` is None (default) → from-scratch: the renderer paints
    RA backgrounds, the brand-mark, the cover orb, and the page-number footer.

    ``template_path`` set → template mode: loads ra_template.pptx, reuses its
    cover for the first slide, appends an ``RA_thankyou`` slide for closing,
    and paints only content on ``RA_blank``/``RA_dark`` for everything else
    (``_TEMPLATE_MODE`` early-returns the chrome helpers)."""
    if template_path:
        pres = Presentation(template_path)
        while len(pres.slides) > 1:
            _remove_slide(pres, len(pres.slides) - 1)
    else:
        pres = Presentation()
        pres.slide_width = Inches(SLIDE_W_IN)
        pres.slide_height = Inches(SLIDE_H_IN)
    pres.core_properties.title = title

    total = len(deck_spec)

    if template_path:
        token = _TEMPLATE_MODE.set(True)
        try:
            content_layout = (_find_layout(pres, _TEMPLATE_CONTENT_LAYOUT)
                              or pres.slide_layouts[6])
            thankyou_layout = _find_layout(pres, _TEMPLATE_THANKYOU_LAYOUT)

            for i, sl in enumerate(deck_spec, start=1):
                sl_type = sl.get("type")
                is_first_cover = (i == 1 and sl_type == "cover")
                is_closing = (sl_type == "closing")

                if is_first_cover:
                    if len(pres.slides) >= 1:
                        slide = pres.slides[0]
                    else:
                        cover_layout = _find_layout(pres, _TEMPLATE_COVER_LAYOUT) \
                                       or content_layout
                        slide = pres.slides.add_slide(cover_layout)
                    _render_template_cover(slide, sl)
                    continue

                if i == 1 and len(pres.slides) >= 1:
                    _remove_slide(pres, 0)

                if is_closing and thankyou_layout is not None:
                    slide = pres.slides.add_slide(thankyou_layout)
                    _render_template_thankyou(slide, sl)
                    continue

                bg = sl.get("bg") or "paper"
                lay_for_slide = _pick_template_layout(
                    pres, sl_type, bg, content_layout,
                )
                slide = pres.slides.add_slide(lay_for_slide)
                _strip_empty_placeholders(slide)
                fill_hex, text_hex, accent_hex = _bg_colours(bg)
                renderer = _RENDERERS.get(sl_type)
                if renderer is None:
                    continue
                try:
                    renderer(slide, sl, i, total, fill_hex, text_hex, accent_hex)
                except Exception as e:  # noqa: BLE001
                    logger.exception("[compose_deck] slide[%d] (%s) render failed",
                                     i - 1, sl_type)
                    raise RuntimeError(f"slide[{i-1}] ({sl_type}): {e}") from e
        finally:
            _TEMPLATE_MODE.reset(token)
    else:
        blank_layout = pres.slide_layouts[6]
        for i, sl in enumerate(deck_spec, start=1):
            slide = pres.slides.add_slide(blank_layout)
            bg = sl.get("bg") or "paper"
            fill_hex, text_hex, accent_hex = _bg_colours(bg)
            renderer = _RENDERERS.get(sl.get("type"))
            if renderer is None:
                continue
            try:
                renderer(slide, sl, i, total, fill_hex, text_hex, accent_hex)
            except Exception as e:  # noqa: BLE001
                logger.exception("[compose_deck] slide[%d] (%s) render failed",
                                 i - 1, sl.get("type"))
                raise RuntimeError(f"slide[{i-1}] ({sl.get('type')}): {e}") from e

    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# HTML rendering (preview) — Resonance Analytics, faithful to the design CSS
# ---------------------------------------------------------------------------


_HTML_HEAD = """<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8" />
<title>__TITLE__</title>
<meta name="viewport" content="width=1920" />
<link rel="preconnect" href="https://fonts.googleapis.com" />
<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin />
<link rel="stylesheet" href="https://fonts.googleapis.com/css2?family=Manrope:wght@300;400;500;600;700;800&family=Playfair+Display:ital,wght@0,400..900;1,400..900&family=JetBrains+Mono:wght@400;500&display=swap" />
<style>
:root {
  --signal-50:#F0F5FF; --signal-100:#DAE9FF; --signal-300:#85B0FF;
  --signal-500:#345FCF; --signal-600:#254BB2; --signal-700:#17368D;
  --signal-800:#0A2163; --signal-900:#051139;
  --ink-100:#EFF2F5; --ink-200:#E1E5E9; --ink-300:#C6CBD0; --ink-400:#9A9FA5;
  --ink-500:#6D7277; --ink-600:#494E52; --ink-700:#2F3338; --ink-900:#050A0F; --ink-950:#010204;
  --gold:#DF9B44; --gold-deep:#B46D10; --teal:#2695AC; --plum:#913F82; --slate:#52657A;
  --paper:#FFFFFF; --cream:#FBF0E4; --border:var(--ink-200);
  --text:var(--ink-900); --text-muted:var(--ink-600); --accent:var(--signal-600);
  --success:#14874E; --danger:#BA2B2E;
  --font-sans:"Manrope",-apple-system,BlinkMacSystemFont,"Segoe UI",sans-serif;
  --font-serif:"Playfair Display","Source Serif 4",Georgia,serif;
  --font-mono:"JetBrains Mono",ui-monospace,"SF Mono",monospace;
}
* { box-sizing: border-box; }
html, body { margin: 0; padding: 0; background: #1a1614; font-family: var(--font-sans);
  -webkit-font-smoothing: antialiased; font-variant-numeric: tabular-nums; }
.slide {
  width: 1920px; height: 1080px; background: var(--paper); color: var(--text);
  padding: 72px 96px 120px; position: relative; overflow: hidden;
  display: flex; flex-direction: column; margin: 24px auto;
  letter-spacing: -0.005em;
}
.slide.bg-signal { background: var(--signal-700); color: #fff; }
.slide.bg-ink    { background: var(--ink-950); color: var(--ink-100); }
.slide.bg-cream  { background: var(--cream); }
.bg-ink .muted, .bg-signal .muted { color: rgba(255,255,255,.72); }
.eyebrow { font-size: 13px; font-weight: 700; letter-spacing: 0.16em;
  text-transform: uppercase; color: var(--accent); display: inline-flex;
  align-items: center; gap: 12px; margin-bottom: 18px; }
.eyebrow::before { content: ""; width: 28px; height: 2px; background: var(--accent); }
.bg-signal .eyebrow, .bg-ink .eyebrow { color: var(--gold); }
.bg-signal .eyebrow::before, .bg-ink .eyebrow::before { background: var(--gold); }
.slide-foot { position: absolute; left: 96px; right: 96px; bottom: 36px;
  display: flex; justify-content: space-between; font-size: 13px;
  color: var(--ink-500); letter-spacing: 0.06em; }
.slide-foot .num { font-family: var(--font-mono); font-weight: 600; color: var(--text); }
.bg-signal .slide-foot, .bg-ink .slide-foot { color: rgba(255,255,255,.6); }
.bg-signal .slide-foot .num, .bg-ink .slide-foot .num { color: #fff; }
.t-title-hero { font-family: var(--font-serif); font-size: 96px; font-weight: 400;
  line-height: 1.0; letter-spacing: -0.025em; }
.t-title-l    { font-family: var(--font-serif); font-size: 64px; font-weight: 400;
  line-height: 1.05; letter-spacing: -0.02em; }
.t-title-m    { font-size: 40px; font-weight: 600; line-height: 1.15; letter-spacing: -0.015em; }
.t-body-l     { font-size: 22px; line-height: 1.5; }
.t-body       { font-size: 18px; line-height: 1.5; }
.muted        { color: var(--text-muted); }
.lay-2 { display: grid; grid-template-columns: 1fr 1fr; gap: 64px; }
.card { background: var(--paper); border: 1px solid var(--border); border-radius: 8px; padding: 32px; }
.bg-ink .card { background: var(--ink-900); border-color: var(--ink-700); }
.kpi { display: flex; flex-direction: column; gap: 4px; position: relative; }
.kpi::before { content:""; position:absolute; top:0; left:0; width:32px; height:3px; background:var(--accent); }
.kpi .label { font-size: 12px; font-weight: 700; letter-spacing: 0.12em;
  text-transform: uppercase; color: var(--ink-600); margin-top: 14px; }
.kpi .value { font-family: var(--font-serif); font-size: 64px; font-weight: 400; color: var(--text); letter-spacing: -0.02em; }
.kpi .unit  { font-size: 24px; font-weight: 600; color: var(--ink-500); margin-left: 6px; font-family: var(--font-sans); }
.kpi .delta { font-size: 16px; font-weight: 700; margin-top: 6px; }
.delta.up   { color: var(--success); }
.delta.down { color: var(--danger); }
.delta.flat { color: var(--ink-500); }
.bullets li { margin-bottom: 14px; }
.bullets li::marker { color: var(--accent); }
table { width: 100%; border-collapse: collapse; }
th, td { padding: 12px 14px; text-align: left; font-size: 15px; }
th { background: var(--signal-600); color: #fff; font-weight: 700;
  font-size: 11px; letter-spacing: 0.06em; text-transform: uppercase; }
tr:nth-child(even) td { background: var(--ink-100); }
.num-col { text-align: right; }
svg.chart { width: 100%; height: 100%; }
.commentary p { margin: 0 0 14px; font-size: 14px; line-height: 1.5; }
.commentary p strong { color: var(--accent); }
.timeline { display: grid; gap: 24px; }
.timeline .step { border-top: 3px solid var(--accent); padding-top: 12px; }
.timeline .step .label { color: var(--accent); font-weight: 700;
  font-size: 11px; letter-spacing: 0.12em; text-transform: uppercase; }
.timeline .step .heading { font-size: 22px; font-weight: 600; margin: 8px 0 12px; }
.timeline .step .body { color: var(--ink-600); font-size: 14px; line-height: 1.5; }
blockquote { font-family: var(--font-serif); margin: 0; padding: 0 64px; font-size: 40px;
  font-style: italic; line-height: 1.3; font-weight: 400; }
.attribution { color: var(--ink-600); font-weight: 700; font-size: 18px; margin-top: 24px; padding: 0 64px; }
.caption { font-size: 14px; color: var(--ink-500); margin-top: 12px; }
.bg-ink .caption, .bg-signal .caption { color: rgba(255,255,255,.6); }
.stat-value { font-family: var(--font-serif); font-size: 180px; font-weight: 400;
  color: var(--accent); line-height: 1; letter-spacing: -0.03em; }
.stat-unit  { font-size: 44px; font-weight: 600; color: var(--ink-600); margin-left: 12px; font-family: var(--font-sans); }
/* Brand-mark glyph */
.brand-mark { position: absolute; top: 56px; left: 96px; display: inline-flex;
  align-items: center; gap: 12px; font-size: 18px; font-weight: 700; color: var(--text); }
.bg-ink .brand-mark, .bg-signal .brand-mark { color: #fff; }
.brand-mark .mk { width: 28px; height: 28px; border-radius: 7px; background: var(--accent);
  position: relative; display: inline-block; }
.bg-ink .brand-mark .mk, .bg-signal .brand-mark .mk { background: #fff; }
.brand-mark .mk::before { content:""; position:absolute; inset:9px 9px auto auto; width:9px; height:2px; background:var(--paper); }
.brand-mark .mk::after { content:""; position:absolute; inset:9px auto auto 9px; width:2px; height:11px; background:var(--paper); }
.bg-ink .brand-mark .mk::before, .bg-ink .brand-mark .mk::after,
.bg-signal .brand-mark .mk::before, .bg-signal .brand-mark .mk::after { background: var(--signal-700); }
/* Cover poster (the design's signature dark cobalt orb) */
.cover-poster { position: absolute; inset: 0; overflow: hidden;
  background:
    radial-gradient(120% 90% at 110% -10%, var(--signal-700), transparent 60%),
    radial-gradient(80% 60% at -10% 110%, oklch(0.32 0.08 35 / .55), transparent 60%),
    linear-gradient(180deg, var(--ink-950), oklch(0.10 0.025 260)); }
.cover-poster::before { content:""; position:absolute; right:-240px; top:-220px;
  width:920px; height:920px; border-radius:50%;
  background:
    radial-gradient(circle at 30% 30%, oklch(0.78 0.13 70 / .55), transparent 55%),
    radial-gradient(circle at 70% 70%, var(--signal-500), var(--signal-800) 60%, var(--signal-900));
  box-shadow: 0 60px 120px -20px rgba(0,0,0,.5); }
.cover-poster::after { content:""; position:absolute; left:-180px; bottom:-260px;
  width:620px; height:620px; border-radius:50%;
  background: radial-gradient(circle at 35% 35%, oklch(0.72 0.13 70 / .35), transparent 65%); filter: blur(4px); }
</style>
</head>
<body>
"""

_HTML_FOOT = "</body></html>\n"


def _esc(s: str) -> str:
    return html.escape(str(s or ""), quote=True)


def _bg_class(bg: str) -> str:
    bg = (bg or "").lower()
    if bg in ("signal", "red"):
        return " bg-signal"
    if bg in ("ink", "charcoal", "dark"):
        return " bg-ink"
    if bg == "cream":
        return " bg-cream"
    return ""


def _foot(idx: int, total: int, note: str = "") -> str:
    return (
        f'<div class="slide-foot">'
        f'<span>{_esc(note)}</span>'
        f'<span class="num">{idx:02d} / {total:02d}</span>'
        f'</div>'
    )


def _brandmark_html() -> str:
    return '<div class="brand-mark"><span class="mk"></span> Resonance Analytics</div>'


def _svg_chart(sl: dict) -> str:
    kind = sl.get("chart", "column")
    cats = list(map(str, sl.get("categories", [])))
    series = sl.get("series", [])
    if not cats or not series:
        return '<div style="font-size:14px;color:#999;">[chart placeholder]</div>'

    W, H = 1680, 600
    PAD_L, PAD_R, PAD_T, PAD_B = 80, 40, 40, 80
    plot_w, plot_h = W - PAD_L - PAD_R, H - PAD_T - PAD_B

    if kind == "pie":
        vals = [float(v) for v in series[0].get("values", [])]
        total = sum(vals) or 1.0
        cx, cy, r = W // 2, H // 2 + 10, min(plot_w, plot_h) // 2 - 20
        out = [f'<svg class="chart" viewBox="0 0 {W} {H}" xmlns="http://www.w3.org/2000/svg">']
        ang = -90.0
        from math import cos, radians, sin
        for i, v in enumerate(vals):
            frac = v / total
            sweep = frac * 360
            a1 = radians(ang)
            a2 = radians(ang + sweep)
            x1 = cx + r * cos(a1); y1 = cy + r * sin(a1)
            x2 = cx + r * cos(a2); y2 = cy + r * sin(a2)
            large = 1 if sweep > 180 else 0
            color = CHART_PALETTE[i % len(CHART_PALETTE)]
            out.append(
                f'<path d="M{cx},{cy} L{x1:.1f},{y1:.1f} A{r},{r} 0 {large} 1 {x2:.1f},{y2:.1f} Z" '
                f'fill="#{color}"/>'
            )
            mid = radians(ang + sweep / 2)
            lx = cx + (r + 30) * cos(mid)
            ly = cy + (r + 30) * sin(mid)
            out.append(
                f'<text x="{lx:.0f}" y="{ly:.0f}" font-family="Manrope,sans-serif" font-size="16" '
                f'fill="#050A0F" text-anchor="middle">{_esc(cats[i]) if i < len(cats) else ""} '
                f'{int(round(frac * 100))}%</text>'
            )
            ang += sweep
        out.append("</svg>")
        return "".join(out)

    all_vals = [float(v) for s in series for v in s.get("values", [])]
    vmax = max(all_vals) if all_vals else 1.0
    vmin = min(0.0, min(all_vals) if all_vals else 0.0)
    rng = (vmax - vmin) or 1.0

    out = [f'<svg class="chart" viewBox="0 0 {W} {H}" xmlns="http://www.w3.org/2000/svg">']
    out.append(f'<line x1="{PAD_L}" y1="{PAD_T + plot_h}" x2="{W - PAD_R}" y2="{PAD_T + plot_h}" '
               f'stroke="#9A9FA5" stroke-width="1"/>')
    out.append(f'<line x1="{PAD_L}" y1="{PAD_T}" x2="{PAD_L}" y2="{PAD_T + plot_h}" '
               f'stroke="#9A9FA5" stroke-width="1"/>')

    value_axis_is_x = (kind == "bar")
    n_ticks = 4
    tick_vals = [vmin + (rng * i / n_ticks) for i in range(n_ticks + 1)]

    def _fmt_tick(v: float) -> str:
        av = abs(v)
        if av >= 1e9:
            return f"{v/1e9:.1f}B"
        if av >= 1e6:
            return f"{v/1e6:.1f}M"
        if av >= 1e3:
            return f"{v/1e3:.1f}K"
        if av == 0:
            return "0"
        if av < 1:
            return f"{v:.2g}"
        return f"{v:.4g}"

    if value_axis_is_x:
        for i, tv in enumerate(tick_vals):
            xg = PAD_L + (i / n_ticks) * plot_w
            if i not in (0, n_ticks):
                out.append(
                    f'<line x1="{xg:.1f}" y1="{PAD_T}" x2="{xg:.1f}" '
                    f'y2="{PAD_T + plot_h}" stroke="#E1E5E9" '
                    f'stroke-width="1" stroke-dasharray="3 4"/>'
                )
            out.append(
                f'<text x="{xg:.1f}" y="{PAD_T + plot_h + 42:.1f}" '
                f'font-family="Manrope,sans-serif" font-size="12" fill="#6D7277" '
                f'text-anchor="middle">{_fmt_tick(tv)}</text>'
            )
    else:
        for i, tv in enumerate(tick_vals):
            yg = PAD_T + plot_h - (i / n_ticks) * plot_h
            if i not in (0, n_ticks):
                out.append(
                    f'<line x1="{PAD_L}" y1="{yg:.1f}" x2="{W - PAD_R}" '
                    f'y2="{yg:.1f}" stroke="#E1E5E9" stroke-width="1" '
                    f'stroke-dasharray="3 4"/>'
                )
            out.append(
                f'<text x="{PAD_L - 10:.1f}" y="{yg + 4:.1f}" '
                f'font-family="Manrope,sans-serif" font-size="12" fill="#6D7277" '
                f'text-anchor="end">{_fmt_tick(tv)}</text>'
            )

    n = len(cats)
    if kind in ("column", "bar"):
        band = (plot_w if kind == "column" else plot_h) / n
        n_series = len(series)
        bar_thick = (band * 0.7) / max(1, n_series)
        for si, s in enumerate(series):
            color = CHART_PALETTE[si % len(CHART_PALETTE)]
            vals = [float(v) for v in s.get("values", [])]
            for ci in range(n):
                v = vals[ci] if ci < len(vals) else 0.0
                if kind == "column":
                    h = (v - vmin) / rng * plot_h
                    x = PAD_L + ci * band + (band * 0.15) + si * bar_thick
                    y = PAD_T + plot_h - h
                    out.append(f'<rect x="{x:.1f}" y="{y:.1f}" width="{bar_thick:.1f}" height="{h:.1f}" fill="#{color}" rx="2"/>')
                    if sl.get("show_values", True):
                        out.append(f'<text x="{x + bar_thick / 2:.1f}" y="{y - 6:.1f}" '
                                   f'font-family="Manrope,sans-serif" font-size="13" fill="#050A0F" '
                                   f'text-anchor="middle">{v:.4g}</text>')
                else:
                    w = (v - vmin) / rng * plot_w
                    y = PAD_T + ci * band + (band * 0.15) + si * bar_thick
                    x = PAD_L
                    out.append(f'<rect x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{bar_thick:.1f}" fill="#{color}" rx="2"/>')
                    if sl.get("show_values", True):
                        out.append(f'<text x="{x + w + 6:.1f}" y="{y + bar_thick / 2 + 5:.1f}" '
                                   f'font-family="Manrope,sans-serif" font-size="13" fill="#050A0F">{v:.4g}</text>')
        if kind == "column":
            for ci, c in enumerate(cats):
                xc = PAD_L + ci * band + band / 2
                out.append(f'<text x="{xc:.1f}" y="{PAD_T + plot_h + 24:.1f}" '
                           f'font-family="Manrope,sans-serif" font-size="14" fill="#494E52" '
                           f'text-anchor="middle">{_esc(c)}</text>')
        else:
            for ci, c in enumerate(cats):
                yc = PAD_T + ci * band + band / 2 + 5
                out.append(f'<text x="{PAD_L - 8:.1f}" y="{yc:.1f}" '
                           f'font-family="Manrope,sans-serif" font-size="14" fill="#494E52" '
                           f'text-anchor="end">{_esc(c)}</text>')
    elif kind == "line":
        for si, s in enumerate(series):
            color = CHART_PALETTE[si % len(CHART_PALETTE)]
            vals = [float(v) for v in s.get("values", [])]
            pts = []
            for ci in range(n):
                v = vals[ci] if ci < len(vals) else 0.0
                x = PAD_L + ci * (plot_w / max(1, n - 1))
                y = PAD_T + plot_h - (v - vmin) / rng * plot_h
                pts.append((x, y))
            d = "M " + " L ".join(f"{x:.1f},{y:.1f}" for x, y in pts)
            out.append(f'<path d="{d}" fill="none" stroke="#{color}" stroke-width="3"/>')
            for x, y in pts:
                out.append(f'<circle cx="{x:.1f}" cy="{y:.1f}" r="4.5" fill="#{color}"/>')
        for ci, c in enumerate(cats):
            xc = PAD_L + ci * (plot_w / max(1, n - 1))
            out.append(f'<text x="{xc:.1f}" y="{PAD_T + plot_h + 24:.1f}" '
                       f'font-family="Manrope,sans-serif" font-size="14" fill="#494E52" '
                       f'text-anchor="middle">{_esc(c)}</text>')

    out.append("</svg>")
    return "".join(out)


def _slide_head_html(sl: dict) -> str:
    parts = []
    if sl.get("eyebrow"):
        parts.append(f'<div class="eyebrow">{_esc(sl["eyebrow"])}</div>')
    parts.append(f'<h2 class="t-title-m" style="margin: 0 0 8px;">{_esc(sl.get("title", ""))}</h2>')
    if sl.get("lede"):
        parts.append(f'<p class="t-body muted" style="margin: 0 0 24px;">{_esc(sl["lede"])}</p>')
    return "".join(parts)


def _render_slide_html(sl: dict, idx: int, total: int) -> str:
    t = sl.get("type")
    bg = sl.get("bg") or ("ink" if t in ("cover", "closing") else "paper")
    cls = f"slide{_bg_class(bg)}"
    head = _slide_head_html(sl)
    foot = _foot(idx, total, sl.get("footer_note", ""))
    on_dark = bg in ("signal", "ink", "red", "charcoal", "dark")

    if t == "cover":
        poster = '<div class="cover-poster"></div>' if on_dark and sl.get("motif", True) else ""
        return (
            f'<section class="{cls}">'
            f'{poster}'
            f'{_brandmark_html()}'
            f'<div style="margin: auto 0; position: relative;">'
            f'<div class="eyebrow">{_esc(sl.get("eyebrow", ""))}</div>'
            f'<h1 class="t-title-hero" style="margin: 18px 0 18px;">{_esc(sl.get("title", ""))}</h1>'
            f'<p class="t-body-l muted" style="max-width: 60ch;">{_esc(sl.get("lede", ""))}</p>'
            f'</div>'
            f'{foot}'
            f'</section>'
        )

    if t == "section_divider":
        return (
            f'<section class="{cls}">'
            f'{_brandmark_html()}'
            f'<div style="margin: auto 0;">'
            f'<div class="eyebrow">{_esc(sl.get("eyebrow", ""))}</div>'
            f'<h2 class="t-title-l" style="margin: 18px 0;">{_esc(sl.get("title", ""))}</h2>'
            f'<p class="t-body-l muted" style="max-width: 70ch;">{_esc(sl.get("lede", ""))}</p>'
            f'</div>'
            f'{foot}'
            f'</section>'
        )

    if t == "stat_callout":
        delta = ""
        if sl.get("delta"):
            dd = sl.get("delta_dir", "flat")
            arrow = "▲ " if dd == "up" else ("▼ " if dd == "down" else "")
            delta = f'<div class="delta {dd}" style="font-size:28px; margin-top:16px;">{arrow}{_esc(sl["delta"])}</div>'
        cap = f'<div class="caption">{_esc(sl["caption"])}</div>' if sl.get("caption") else ""
        return (
            f'<section class="{cls}">'
            f'{head}'
            f'<div style="margin: 60px 0 auto;">'
            f'<span class="stat-value">{_esc(sl.get("value", ""))}</span>'
            f'<span class="stat-unit">{_esc(sl.get("unit", ""))}</span>'
            f'{delta}'
            f'</div>'
            f'{cap}'
            f'{foot}'
            f'</section>'
        )

    if t == "kpi_grid":
        kpis = sl.get("kpis", [])[:6]
        cells = []
        for k in kpis:
            dd = k.get("delta_dir", "flat")
            arrow = "▲ " if dd == "up" else ("▼ " if dd == "down" else "")
            delta = f'<div class="delta {dd}">{arrow}{_esc(k.get("delta", ""))}</div>' if k.get("delta") else ""
            unit_span = f'<span class="unit">{_esc(k.get("unit", ""))}</span>' if k.get("unit") else ""
            cells.append(
                f'<div class="card kpi">'
                f'<div class="label">{_esc(k.get("label", ""))}</div>'
                f'<div class="value">{_esc(k.get("value", ""))}{unit_span}</div>'
                f'{delta}</div>'
            )
        cols = "1fr 1fr 1fr 1fr" if len(kpis) <= 4 else "1fr 1fr 1fr"
        cap = f'<div class="caption" style="margin-top:24px;">{_esc(sl["caption"])}</div>' if sl.get("caption") else ""
        return (
            f'<section class="{cls}">'
            f'{head}'
            f'<div style="display:grid; grid-template-columns:{cols}; gap:24px; flex:1;">{"".join(cells)}</div>'
            f'{cap}'
            f'{foot}'
            f'</section>'
        )

    if t == "bullets":
        items = "".join(f"<li>{_esc(b)}</li>" for b in sl.get("bullets", []))
        return (
            f'<section class="{cls}">'
            f'{head}'
            f'<ul class="bullets t-body-l" style="margin-top:32px;">{items}</ul>'
            f'{foot}'
            f'</section>'
        )

    if t == "two_column":
        def col(side: dict) -> str:
            items = "".join(f"<li>{_esc(b)}</li>" for b in (side.get("bullets") or []))
            return (
                f'<div>'
                f'<h3 class="t-title-m" style="font-size:28px; margin: 0 0 16px;">{_esc(side.get("heading", ""))}</h3>'
                f'<ul class="bullets t-body" style="padding-left:24px;">{items}</ul>'
                f'</div>'
            )
        return (
            f'<section class="{cls}">'
            f'{head}'
            f'<div class="lay-2" style="margin-top:32px; flex:1;">'
            f'{col(sl.get("left", {}) or {})}'
            f'{col(sl.get("right", {}) or {})}'
            f'</div>'
            f'{foot}'
            f'</section>'
        )

    if t == "chart":
        cap = f'<div class="caption" style="margin-top:16px;">{_esc(sl["caption"])}</div>' if sl.get("caption") else ""
        return (
            f'<section class="{cls}">'
            f'{head}'
            f'<div class="card" style="flex:1; margin-top:24px; padding:32px;">{_svg_chart(sl)}</div>'
            f'{cap}'
            f'{foot}'
            f'</section>'
        )

    if t == "chart_commentary":
        items = _coerce_commentary(sl.get("commentary"))[:5]
        c_html = "".join(
            f'<p><strong>{_esc(it.get("label", ""))}</strong> {_esc(it.get("text", ""))}</p>'
            for it in items if isinstance(it, dict)
        )
        return (
            f'<section class="{cls}">'
            f'{head}'
            f'<div style="display:grid; grid-template-columns: 1.6fr 1fr; gap:40px; flex:1; margin-top:24px;">'
            f'<div class="card" style="padding:32px;">{_svg_chart(sl)}</div>'
            f'<div class="commentary">{c_html}</div>'
            f'</div>'
            f'{foot}'
            f'</section>'
        )

    if t == "table":
        headers = sl.get("headers", [])
        rows = sl.get("rows", [])
        num_idx = {
            i for i, h in enumerate(headers)
            if any(tok in str(h).lower() for tok in ("%", "thb", "usd", "eur", "gbp", "(m)", "delta", "δ"))
        }
        thead = "<tr>" + "".join(f"<th>{_esc(h)}</th>" for h in headers) + "</tr>"
        body_rows = []
        for row in rows:
            cells = []
            for ci, v in enumerate(row):
                cls_td = ' class="num-col"' if ci in num_idx else ""
                cells.append(f"<td{cls_td}>{_esc(v)}</td>")
            body_rows.append("<tr>" + "".join(cells) + "</tr>")
        cap = f'<div class="caption" style="margin-top:16px;">{_esc(sl["caption"])}</div>' if sl.get("caption") else ""
        return (
            f'<section class="{cls}">'
            f'{head}'
            f'<div style="flex:1; margin-top:24px;"><table>{thead}{"".join(body_rows)}</table></div>'
            f'{cap}'
            f'{foot}'
            f'</section>'
        )

    if t == "quote":
        return (
            f'<section class="{cls}">'
            f'<div style="margin:auto 0; text-align:center;">'
            f'<div class="eyebrow" style="justify-content:center;">{_esc(sl.get("eyebrow", ""))}</div>'
            f'<blockquote>&ldquo;{_esc(sl.get("quote", ""))}&rdquo;</blockquote>'
            f'<div class="attribution">— {_esc(sl.get("attribution", ""))}</div>'
            f'</div>'
            f'<div class="caption" style="text-align:center; margin-bottom:80px;">{_esc(sl.get("caption", ""))}</div>'
            f'{foot}'
            f'</section>'
        )

    if t == "data_note":
        return (
            f'<section class="{cls}">'
            f'{head}'
            f'<p class="t-body-l muted" style="max-width: 70ch; margin-top:20px;">'
            f'Requested: <strong>{_esc(sl.get("requested", "—"))}</strong>. '
            f'Available data: <strong>{_esc(sl.get("available", "—"))}</strong>. '
            f'{_esc(sl.get("reason", ""))}'
            f'</p>'
            f'{foot}'
            f'</section>'
        )

    if t == "timeline":
        steps = sl.get("steps", [])[:5]
        cols = f"repeat({len(steps)}, 1fr)" if steps else "1fr"
        items = "".join(
            f'<div class="step"><div class="label">{_esc(s.get("label", ""))}</div>'
            f'<div class="heading">{_esc(s.get("heading", ""))}</div>'
            f'<div class="body">{_esc(s.get("body", ""))}</div></div>'
            for s in steps
        )
        return (
            f'<section class="{cls}">'
            f'{head}'
            f'<div class="timeline" style="grid-template-columns:{cols}; flex:1; margin-top:32px;">{items}</div>'
            f'{foot}'
            f'</section>'
        )

    if t == "closing":
        poster = '<div class="cover-poster"></div>' if on_dark and sl.get("motif", True) else ""
        bl = ""
        if sl.get("bullets"):
            lis = "".join(f"<li>{_esc(b)}</li>" for b in sl["bullets"])
            bl = f'<ul class="bullets t-body muted" style="margin-top:20px;max-width:70ch;">{lis}</ul>'
        return (
            f'<section class="{cls}">'
            f'{poster}'
            f'{_brandmark_html()}'
            f'<div style="margin: auto 0; position: relative;">'
            f'<div class="eyebrow">{_esc(sl.get("eyebrow", "") or "NEXT STEPS")}</div>'
            f'<h2 class="t-title-l" style="margin: 18px 0;">{_esc(sl.get("title", ""))}</h2>'
            f'<p class="t-body-l muted" style="max-width: 60ch;">{_esc(sl.get("lede", ""))}</p>'
            f'{bl}'
            f'</div>'
            f'{foot}'
            f'</section>'
        )

    return f'<section class="{cls}"><h2>{_esc(sl.get("type", ""))}</h2></section>'


def _render_html(title: str, deck_spec: list) -> bytes:
    total = len(deck_spec)
    body = "".join(_render_slide_html(sl, i + 1, total) for i, sl in enumerate(deck_spec))
    out = _HTML_HEAD.replace("__TITLE__", _esc(title)) + body + _HTML_FOOT
    return out.encode("utf-8")


# ---------------------------------------------------------------------------
# Tool factory
# ---------------------------------------------------------------------------


def build_compose_deck_tool(
    *,
    workspace_client: Any,
    variable_store_cls: Any = None,
    app_url: str | None = None,
    template_path: str | None = None,
):
    """Build the ``compose_deck`` tool (Resonance Analytics brand).

    Args:
        workspace_client: ``databricks.sdk.WorkspaceClient`` for Files API.
        variable_store_cls: Reserved (unused; kept for factory parity with
            compose_document / compose_infographic).
        app_url: Chat-app base URL for ``/api/decks/<id>`` proxy URLs.
            Defaults to env ``APP_URL``.
        template_path: Optional path to a base ``ra_template.pptx``. When set,
            brand chrome (cover, footer, layouts) is taken from the template
            instead of drawn by the renderer. Resolution order: explicit
            ``template_path`` → env ``COMPOSE_DECK_TEMPLATE`` → bundled
            ``skills/compose-pptx/templates/ra_template.pptx`` (absent by
            default → from-scratch mode). Set the env var to ``""`` to force
            from-scratch mode and skip the bundled default.
    """
    _ = variable_store_cls
    base_app_url = (app_url or os.environ.get("APP_URL", "")).rstrip("/")
    env_tmpl = os.environ.get("COMPOSE_DECK_TEMPLATE")
    if template_path:
        base_template_path = template_path
    elif env_tmpl is not None:
        base_template_path = env_tmpl or None
    elif _DEFAULT_TEMPLATE_PATH.exists():
        base_template_path = str(_DEFAULT_TEMPLATE_PATH)
    else:
        base_template_path = None
    if base_template_path and not os.path.exists(base_template_path):
        logger.warning("[compose_deck] template_path does not exist: %s — "
                       "falling back to from-scratch mode", base_template_path)
        base_template_path = None
    if base_template_path:
        logger.warning("[compose_deck] template mode ON, template=%s",
                       base_template_path)

    @tool
    def compose_deck(
        title: str,
        deck_spec: list,
        config: RunnableConfig = None,
        store: Annotated[Any, InjectedStore()] = None,
    ) -> str:
        """Compose a Resonance Analytics-branded deck → editable PPTX (native charts) + HTML preview.

        Use whenever the user asks for slides, a deck, a presentation, or a
        .pptx briefing. Author a ``deck_spec`` JSON list — one dict per slide
        with ``type`` ∈ {cover, section_divider, stat_callout, kpi_grid,
        bullets, two_column, chart, chart_commentary, table, quote, data_note,
        timeline, closing}. Chart slides take ``chart`` ∈
        {"column","bar","line","pie"} (a STRING) with the data in the slide's
        own ``categories`` + ``series`` keys. FIRST call
        ``find_skill("compose a presentation deck")`` for the slide-spec
        fields, deck-structure guide, and brand rules.

        Args:
            title: Deck title — on-volume filename + cache key (re-rendering
                the same (title, spec) overwrites).
            deck_spec: List of slide dicts, each with ``type``.

        Returns:
            JSON {status, document_id, title, slide_count, preview_url,
            pptx_url, ...} or error{code,message}. Surface ``preview_url`` +
            ``pptx_url`` as links; never the Volumes path.
        """
        _ = config, store

        if not title or not (isinstance(title, str) and title.strip()):
            return json.dumps(_compact_error("empty_title", "title is required"))
        title = title.strip()
        if deck_spec is None:
            return json.dumps(_compact_error("empty_spec", "deck_spec is required"))

        ok, msg = _validate_spec(deck_spec)
        if not ok:
            code = "invalid_chart_type" if "chart=" in msg else (
                "invalid_slide_type" if "type=" in msg else
                "no_slides" if "empty" in msg else
                "render_failed"
            )
            return json.dumps(_compact_error(code, msg))

        try:
            pptx_bytes = _render_pptx(title, deck_spec,
                                      template_path=base_template_path)
        except Exception as e:  # noqa: BLE001
            logger.exception("[compose_deck] PPTX render failed")
            return json.dumps(_compact_error("render_failed",
                              f"pptx render: {str(e)[:300]}"))

        try:
            html_bytes = _render_html(title, deck_spec)
        except Exception as e:  # noqa: BLE001
            logger.exception("[compose_deck] HTML render failed")
            return json.dumps(_compact_error("render_failed",
                              f"html render: {str(e)[:300]}"))

        doc_id = _document_id(title, deck_spec)
        slug = _slug(title)
        base = f"{_VOLUME_ROOT}/{doc_id}__{slug}"

        try:
            workspace_client.files.create_directory(_VOLUME_ROOT)
        except Exception:  # noqa: BLE001
            pass

        uploads = [
            (f"{base}.pptx", pptx_bytes),
            (f"{base}.html", html_bytes),
            (f"{base}.json", json.dumps({"title": title, "deck_spec": deck_spec},
                                        separators=(",", ":")).encode("utf-8")),
        ]
        try:
            for path, blob in uploads:
                workspace_client.files.upload(path, io.BytesIO(blob), overwrite=True)
        except Exception as e:  # noqa: BLE001
            logger.exception("[compose_deck] Files API upload failed")
            return json.dumps(_compact_error(
                "upload_failed",
                f"Files API upload to {_VOLUME_ROOT} failed: {str(e)[:200]}",
            ))

        out: dict[str, Any] = {
            "status": "ok",
            "document_id": doc_id,
            "title": title,
            "slide_count": len(deck_spec),
            "volumes_path": f"{base}.{{pptx,html,json}}",
            "size_bytes": {"pptx": len(pptx_bytes), "html": len(html_bytes)},
        }
        if base_app_url:
            out["preview_url"] = f"{base_app_url}/api/decks/{doc_id}"
            out["pptx_url"]    = f"{base_app_url}/api/decks/{doc_id}.pptx"
        return json.dumps(out, separators=(",", ":"))

    return compose_deck
