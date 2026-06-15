"""Compose office-suite documents (pptx/docx/xlsx/csv/pdf) and upload to UC Volumes.

Sibling to ``compose_infographic`` — same I/O shape, same upload mechanism, but
for binary office-format artifacts instead of HTML. Filed 2026-05-15 after
trace ``tr-83d056e1a485c06f8b19df47db7f08e3`` proved that the agent's
``run_python_code`` improvisation route produces a fake-successful pptx that
lands in the serving container's local tmpfs (NOT Unity Catalog) because the
in-process exec has no UC volume mount.

This tool uses the Databricks Files API (same path as compose_infographic) so
the bytes actually land in UC and can be served back through the app's
``/api/documents/:id`` proxy.

Supported formats and the library each uses:
    pptx  — python-pptx       — list-of-slides with title/methodology/table/bullets
    docx  — python-docx       — outline of headings + paragraphs + tables + bullets
    xlsx  — openpyxl          — one sheet per `sheets[]` entry (or one default)
    csv   — pandas.to_csv     — a single DataFrame ref → CSV bytes
    pdf   — reportlab         — Platypus story: title, paragraphs, tables, bullets

Section schema (same across pptx/docx/pdf; xlsx uses `sheets[]`; csv ignores):
    {"type": "title",     "text": "..."}
    {"type": "heading",   "level": 1|2|3, "text": "..."}
    {"type": "paragraph", "text": "..."}
    {"type": "bullets",   "items": ["...", "..."]}
    {"type": "kv",        "items": {"label": "value", ...}}
    {"type": "table",     "headers": ["A","B"], "rows": [[...], [...]]}

Return shape::

    {"status": "ok",
     "document_id": "document_<format>_<12hex>",
     "format": "pptx|docx|xlsx|csv|pdf",
     "title": "...",
     "url": "/api/documents/<id>",
     "size_bytes": <int>}
"""

from __future__ import annotations

import csv as csv_module
import hashlib
import io
import json
import logging
import re
from io import BytesIO, StringIO
from typing import Annotated, Any, Callable

import pandas as pd
from langchain_core.runnables import RunnableConfig
from langchain_core.tools import tool
from langgraph.prebuilt import InjectedStore

from tools.compact_ref import _compact_error

logger = logging.getLogger(__name__)


_SUPPORTED_FORMATS = ("pptx", "docx", "xlsx", "csv", "pdf")
_VOLUME_ROOT = "/Volumes/workspace/ai_ops/agent_scratch/documents"


def _document_id(fmt: str, title: str) -> str:
    """Stable id of the form ``document_<format>_<12hex>`` keyed on (fmt, title)."""
    h = hashlib.sha256(f"{fmt}|{title}".encode("utf-8")).hexdigest()[:12]
    return f"document_{fmt}_{h}"


def _slug(title: str, fallback: str = "document") -> str:
    """A filesystem-safe slug for the downloaded filename."""
    s = re.sub(r"[^A-Za-z0-9._-]+", "_", title.strip())
    s = re.sub(r"_+", "_", s).strip("_")
    return (s or fallback)[:80]


# ── Format builders ───────────────────────────────────────────────────────

def _build_pptx(*, title: str, sections: list[dict]) -> bytes:
    """One slide per section. `title` becomes the first slide."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide always first
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = ""

    blank_layout = prs.slide_layouts[6]
    title_only_layout = prs.slide_layouts[5]

    for sec in sections:
        kind = sec.get("type")
        if kind == "title":
            # Skip — first slide already covers it.
            continue
        s = prs.slides.add_slide(title_only_layout)
        if kind == "heading":
            s.shapes.title.text = sec.get("text", "")
        elif kind == "paragraph":
            s.shapes.title.text = ""
            tx = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(5.5))
            tx.text_frame.text = sec.get("text", "")
            for p in tx.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(18)
        elif kind == "bullets":
            s.shapes.title.text = sec.get("heading", "")
            tx = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(5.5))
            tf = tx.text_frame
            for i, item in enumerate(sec.get("items") or []):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {item}"
                for r in p.runs:
                    r.font.size = Pt(18)
        elif kind == "kv":
            s.shapes.title.text = sec.get("heading", "")
            tx = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(5.5))
            tf = tx.text_frame
            items = sec.get("items") or {}
            for i, (k, v) in enumerate(items.items()):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"{k}: {v}"
                for r in p.runs:
                    r.font.size = Pt(16)
        elif kind == "table":
            s.shapes.title.text = sec.get("heading", "")
            headers = sec.get("headers") or []
            rows = sec.get("rows") or []
            if not headers:
                continue
            tbl = s.shapes.add_table(
                rows=len(rows) + 1, cols=len(headers),
                left=Inches(0.5), top=Inches(1.4),
                width=Inches(12), height=Inches(5),
            ).table
            for j, h in enumerate(headers):
                tbl.cell(0, j).text = str(h)
            for i, row in enumerate(rows, start=1):
                for j, cell in enumerate(row):
                    if j < len(headers):
                        tbl.cell(i, j).text = str(cell) if cell is not None else ""
        # Unknown types silently skipped — agent can recover by re-emitting.

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_docx(*, title: str, sections: list[dict]) -> bytes:
    """Flowing document — title, then sections rendered inline."""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(title, level=0)

    for sec in sections:
        kind = sec.get("type")
        if kind == "title":
            continue
        if kind == "heading":
            level = int(sec.get("level", 1)) or 1
            doc.add_heading(sec.get("text", ""), level=min(max(level, 1), 4))
        elif kind == "paragraph":
            doc.add_paragraph(sec.get("text", ""))
        elif kind == "bullets":
            heading = sec.get("heading")
            if heading:
                doc.add_heading(heading, level=2)
            for item in sec.get("items") or []:
                doc.add_paragraph(str(item), style="List Bullet")
        elif kind == "kv":
            heading = sec.get("heading")
            if heading:
                doc.add_heading(heading, level=2)
            for k, v in (sec.get("items") or {}).items():
                p = doc.add_paragraph()
                run = p.add_run(f"{k}: ")
                run.bold = True
                p.add_run(str(v))
        elif kind == "table":
            heading = sec.get("heading")
            if heading:
                doc.add_heading(heading, level=2)
            headers = sec.get("headers") or []
            rows = sec.get("rows") or []
            if not headers:
                continue
            t = doc.add_table(rows=len(rows) + 1, cols=len(headers))
            t.style = "Light Grid"
            for j, h in enumerate(headers):
                t.cell(0, j).text = str(h)
            for i, row in enumerate(rows, start=1):
                for j, cell in enumerate(row):
                    if j < len(headers):
                        t.cell(i, j).text = str(cell) if cell is not None else ""

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_xlsx(*, title: str, sheets: list[dict] | None, dataframes: dict[str, pd.DataFrame]) -> bytes:
    """Multi-sheet workbook. ``sheets`` is a list of:
        {"name": "...", "variable_name": "stored_df"}  -- load from VariableStore
        {"name": "...", "headers": [...], "rows": [[...]]}  -- inline data
    Falls back to single sheet from the first dataframe if `sheets` is empty.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    wb.remove(wb.active)

    items = sheets or []
    if not items and dataframes:
        first_name, first_df = next(iter(dataframes.items()))
        items = [{"name": first_name[:31], "variable_name": first_name}]
    if not items:
        items = [{"name": "Sheet1", "headers": ["—"], "rows": [["(no data provided)"]]}]

    for sheet_spec in items:
        sheet_name = (sheet_spec.get("name") or "Sheet")[:31]
        ws = wb.create_sheet(title=sheet_name)
        if "variable_name" in sheet_spec and sheet_spec["variable_name"] in dataframes:
            df = dataframes[sheet_spec["variable_name"]]
            ws.append(list(df.columns))
            for cell in ws[1]:
                cell.font = Font(bold=True)
            for _, row in df.iterrows():
                ws.append([row[c] if pd.notna(row[c]) else None for c in df.columns])
        else:
            headers = sheet_spec.get("headers") or []
            rows = sheet_spec.get("rows") or []
            if headers:
                ws.append([str(h) for h in headers])
                for cell in ws[1]:
                    cell.font = Font(bold=True)
            for row in rows:
                ws.append(list(row))

    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _build_csv(*, dataframes: dict[str, pd.DataFrame], variable_name: str | None) -> bytes:
    """CSV from a single DataFrame ref. Picks the first dataframe if no ref."""
    if variable_name and variable_name in dataframes:
        df = dataframes[variable_name]
    elif dataframes:
        df = next(iter(dataframes.values()))
    else:
        return b"(no data)\n"
    buf = StringIO()
    df.to_csv(buf, index=False, quoting=csv_module.QUOTE_MINIMAL)
    return buf.getvalue().encode("utf-8")


def _build_pdf(*, title: str, sections: list[dict]) -> bytes:
    """ReportLab Platypus story: flowables only, no fancy layout."""
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem,
    )

    styles = getSampleStyleSheet()
    h1 = styles["Heading1"]
    h2 = styles["Heading2"]
    h3 = styles["Heading3"]
    body = styles["BodyText"]
    title_style = ParagraphStyle("title", parent=h1, fontSize=22, spaceAfter=18)
    kv_style = ParagraphStyle("kv", parent=body, fontSize=11, leading=14)

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=LETTER,
        leftMargin=0.7 * inch, rightMargin=0.7 * inch,
        topMargin=0.7 * inch, bottomMargin=0.7 * inch,
    )
    story = [Paragraph(title, title_style)]

    for sec in sections:
        kind = sec.get("type")
        if kind == "title":
            continue
        if kind == "heading":
            level = int(sec.get("level", 1)) or 1
            style = {1: h1, 2: h2}.get(level, h3)
            story.append(Paragraph(sec.get("text", ""), style))
        elif kind == "paragraph":
            story.append(Paragraph(sec.get("text", ""), body))
            story.append(Spacer(1, 6))
        elif kind == "bullets":
            heading = sec.get("heading")
            if heading:
                story.append(Paragraph(heading, h2))
            items = [ListItem(Paragraph(str(it), body)) for it in (sec.get("items") or [])]
            if items:
                story.append(ListFlowable(items, bulletType="bullet"))
                story.append(Spacer(1, 6))
        elif kind == "kv":
            heading = sec.get("heading")
            if heading:
                story.append(Paragraph(heading, h2))
            for k, v in (sec.get("items") or {}).items():
                story.append(Paragraph(f"<b>{k}:</b> {v}", kv_style))
            story.append(Spacer(1, 6))
        elif kind == "table":
            heading = sec.get("heading")
            if heading:
                story.append(Paragraph(heading, h2))
            headers = sec.get("headers") or []
            rows = sec.get("rows") or []
            if headers:
                data = [list(headers)] + [list(r) for r in rows]
                t = Table(data, hAlign="LEFT")
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]))
                story.append(t)
                story.append(Spacer(1, 8))

    doc.build(story)
    return buf.getvalue()


# ── Public factory ────────────────────────────────────────────────────────

def build_compose_document_tool(
    *,
    workspace_client: Any,
    variable_store_cls: Any,
    app_url: str | None = None,
):
    """Build the ``compose_document`` tool.

    Args:
        workspace_client: ``databricks.sdk.WorkspaceClient`` used for Files API upload.
        variable_store_cls: VariableStore class for hydrating DataFrame refs.
        app_url: Base URL of the chat app (e.g. ``https://pudding-chatbot-….aws.databricksapps.com``)
            for emitting absolute proxy URLs. If unset, returns app-relative paths.
    """

    @tool
    def compose_document(
        format: str,
        title: str,
        sections: list[dict] | None = None,
        sheets: list[dict] | None = None,
        variable_names: list[str] | None = None,
        variable_name: str | None = None,
        config: RunnableConfig = None,
        store: Annotated[Any, InjectedStore()] = None,
    ) -> str:
        """Compose a binary office document, upload to UC Volumes, return a download URL.

        Use for ANY binary-file request (Word, Excel, CSV, PDF). NEVER build
        binary files in ``run_python_code`` — that container has no UC Volume
        mount; the bytes are silently lost. For branded presentation decks
        prefer ``compose_deck``; this tool covers everything else.

        Args:
            format: One of "pptx", "docx", "xlsx", "csv", "pdf".
            title: Document title (title slide / H1 / filename slug).
            sections: pptx/docx/pdf content blocks:
                ``{"type":"heading","level":1,"text"}``, ``{"type":"paragraph","text"}``,
                ``{"type":"bullets","heading","items":[..]}``,
                ``{"type":"kv","heading","items":{label:value}}``,
                ``{"type":"table","heading","headers":[..],"rows":[[..]]}``.
            sheets: xlsx only — ``[{"name","variable_name"}]`` or
                ``[{"name","headers","rows"}]``.
            variable_names: Stored DataFrame names to hydrate for table/sheet refs.
            variable_name: csv only — single stored DataFrame ref.

        Returns:
            JSON ``{status, document_id, format, title, url, size_bytes}`` on
            success; compact error payload otherwise. The frontend surfaces it
            as a download card — reference the document by TITLE in prose.
        """
        fmt = (format or "").lower().strip()
        if fmt not in _SUPPORTED_FORMATS:
            return json.dumps(_compact_error(
                "bad_format",
                f"format must be one of {_SUPPORTED_FORMATS}, got {format!r}",
            ))
        if not title or not title.strip():
            return json.dumps(_compact_error("empty_title", "title is required"))
        title = title.strip()
        sections = sections or []

        # Hydrate any referenced DataFrames from VariableStore
        dataframes: dict[str, pd.DataFrame] = {}
        names_to_load = set(variable_names or [])
        if variable_name:
            names_to_load.add(variable_name)
        for sec in sections:
            if sec.get("type") == "table" and sec.get("variable_name"):
                names_to_load.add(sec["variable_name"])
        for sh in (sheets or []):
            if sh.get("variable_name"):
                names_to_load.add(sh["variable_name"])

        if names_to_load and store is not None and variable_store_cls is not None:
            try:
                proxy = variable_store_cls(store, config or {})
                for name in names_to_load:
                    try:
                        df = proxy.get(name)
                        if df is not None:
                            dataframes[name] = df
                    except Exception as e:
                        logger.warning("[compose_document] failed to load %s: %s", name, e)
            except Exception as e:
                logger.warning("[compose_document] variable_store init failed: %s", e)

        # Build bytes
        try:
            if fmt == "pptx":
                payload = _build_pptx(title=title, sections=sections)
            elif fmt == "docx":
                payload = _build_docx(title=title, sections=sections)
            elif fmt == "xlsx":
                payload = _build_xlsx(title=title, sheets=sheets, dataframes=dataframes)
            elif fmt == "csv":
                payload = _build_csv(dataframes=dataframes, variable_name=variable_name)
            elif fmt == "pdf":
                payload = _build_pdf(title=title, sections=sections)
            else:
                # Defensive — guarded above, but keep the path explicit.
                return json.dumps(_compact_error("bad_format", f"unsupported format {fmt!r}"))
        except ImportError as e:
            return json.dumps(_compact_error(
                "missing_dep",
                f"Builder library not installed: {e}. Check orchestrator requirements.txt.",
            ))
        except Exception as e:
            logger.exception("[compose_document] build_%s failed", fmt)
            return json.dumps(_compact_error("build_error", str(e)[:200]))

        # Upload via Files API to UC Volumes
        doc_id = _document_id(fmt, title)
        slug = _slug(title)
        volumes_path = f"{_VOLUME_ROOT}/{doc_id}__{slug}.{fmt}"
        try:
            try:
                workspace_client.files.create_directory(_VOLUME_ROOT)
            except Exception:
                pass
            workspace_client.files.upload(
                volumes_path,
                BytesIO(payload),
                overwrite=True,
            )
        except Exception as e:
            logger.exception("[compose_document] upload failed")
            return json.dumps(_compact_error(
                "upload_failed",
                f"Files API upload to {_VOLUME_ROOT} failed: {str(e)[:200]}",
            ))

        # Emit ONLY the proxied app URL — never the raw Volumes path. Mirrors
        # the compose_infographic guard against LLM URL-fabrication.
        url = (
            f"{app_url.rstrip('/')}/api/documents/{doc_id}"
            if app_url else f"/api/documents/{doc_id}"
        )
        return json.dumps({
            "status": "ok",
            "document_id": doc_id,
            "format": fmt,
            "title": title,
            "url": url,
            "size_bytes": len(payload),
        }, separators=(",", ":"))

    return compose_document
